import os
import pandas as pd
from datetime import datetime, timedelta
from docx import Document
from docx.shared import Pt, Inches, Cm
from docx.enum.text import WD_PARAGRAPH_ALIGNMENT
from docx.enum.table import WD_ALIGN_VERTICAL
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.enum.text import WD_BREAK
import json
import math

import zipfile
from io import BytesIO
import xml.etree.ElementTree as ET

import textwrap, re

# -------------------------------
# CONFIG
# -------------------------------

EXCEL_NAME = "Main Liturgy file.xlsx"
SHEET_NAME = "Data"
LOGO_FILE = "logo.png"
REGIO = "AM"
PREEK_NAME = "Preek.docx"
HMH = "Haleluja"

bible_books_nl = {
    # Oude Testament
    "Genesis": "GEN",
    "Exodus": "EXO",
    "Leviticus": "LEV",
    "Numeri": "NUM",
    "Deuteronomium": "DEU",
    "Jozua": "JOS",
    "Richteren": "JDG",
    "Ruth": "RUT",
    "1 Samuël": "1SA",
    "2 Samuël": "2SA",
    "1 Koningen": "1KI",
    "2 Koningen": "2KI",
    "1 Kronieken": "1CH",
    "2 Kronieken": "2CH",
    "Ezra": "EZR",
    "Nehemia": "NEH",
    "Ester": "EST",
    "Job": "JOB",
    "Psalmen": "PSA",
    "Spreuken": "PRO",
    "Prediker": "ECC",
    "Hooglied": "SNG",
    "Jesaja": "ISA",
    "Jeremia": "JER",
    "Klaagliederen": "LAM",
    "Ezechiël": "EZK",
    "Daniël": "DAN",
    "Hosea": "HOS",
    "Joël": "JOL",
    "Amos": "AMO",
    "Obadja": "OBA",
    "Jona": "JON",
    "Micha": "MIC",
    "Nahum": "NAM",
    "Habakuk": "HAB",
    "Sefanja": "ZEP",
    "Haggai": "HAG",
    "Zacharia": "ZEC",
    "Maleachi": "MAL",

    # Nieuwe Testament
    "Mattheüs": "MAT",
    "Marcus": "MRK",
    "Lucas": "LUK",
    "Johannes": "JHN",
    "Handelingen": "ACT",
    "Romeinen": "ROM",
    "1 Korintiërs": "1CO",
    "2 Korintiërs": "2CO",
    "Galaten": "GAL",
    "Efeziërs": "EPH",
    "Filippenzen": "PHP",
    "Kolossenzen": "COL",
    "1 Thessalonicenzen": "1TH",
    "2 Thessalonicenzen": "2TH",
    "1 Timoteüs": "1TI",
    "2 Timoteüs": "2TI",
    "Titus": "TIT",
    "Filemon": "PHM",
    "Hebreeën": "HEB",
    "Jakobus": "JAS",   
    "1 Petrus": "1PE",
    "2 Petrus": "2PE",
    "1 Johannes": "1JN",
    "2 Johannes": "2JN",
    "3 Johannes": "3JN",
    "Judas": "JUD",
    "Openbaring": "REV"
}

# Mapping Dutch Bible book names to Indonesian names
dutch_to_indonesian_bible = {
    # Oude Testament
    "Genesis": "Kejadian",
    "Exodus": "Keluaran",
    "Leviticus": "Imamat",
    "Numeri": "Bilangan",
    "Deuteronomium": "Ulangan",
    "Jozua": "Yosua",
    "Richteren": "Hakim-hakim",
    "Ruth": "Rut",
    "1 Samuël": "1 Samuel",
    "2 Samuël": "2 Samuel",
    "1 Koningen": "1 Raja-raja",
    "2 Koningen": "2 Raja-raja",
    "1 Kronieken": "1 Tawarikh",
    "2 Kronieken": "2 Tawarikh",
    "Ezra": "Ezra",
    "Nehemia": "Nehemia",
    "Ester": "Ester",
    "Job": "Ayub",
    "Psalmen": "Mazmur",
    "Spreuken": "Amsal",
    "Prediker": "Pengkhotbah",
    "Hooglied": "Kidung Agung",
    "Jesaja": "Yesaya",
    "Jeremia": "Yeremia",
    "Klaagliederen": "Ratapan",
    "Ezechiël": "Yehezkiel",
    "Daniël": "Daniel",
    "Hosea": "Hosea",
    "Joël": "Yoel",
    "Amos": "Amos",
    "Obadja": "Obaja",
    "Jona": "Yunus",
    "Micha": "Mikha",
    "Nahum": "Nahum",
    "Habakuk": "Habakuk",
    "Sefanja": "Zefanya",
    "Haggaï": "Hagai",
    "Zacharia": "Zakharia",
    "Maleachi": "Maleakhi",

    # Nieuwe Testament
    "Mattheüs": "Matius",
    "Marcus": "Markus",
    "Lucas": "Lukas",
    "Johannes": "Yohanes",
    "Handelingen": "Kisah Para Rasul",
    "Romeinen": "Roma",
    "1 Korintiërs": "1 Korintus",
    "2 Korintiërs": "2 Korintus",
    "Galaten": "Galatia",
    "Efeziërs": "Efesus",
    "Filippenzen": "Filipi",
    "Kolossenzen": "Kolose",
    "1 Thessalonicenzen": "1 Tesalonika",
    "2 Thessalonicenzen": "2 Tesalonika",
    "1 Timoteüs": "1 Timotius",
    "2 Timoteüs": "2 Timotius",
    "Titus": "Titus",
    "Filemon": "Filemon",
    "Hebreeën": "Ibrani",
    "Jakobus": "Yakobus",
    "1 Petrus": "1 Petrus",
    "2 Petrus": "2 Petrus",
    "1 Johannes": "1 Yohanes",
    "2 Johannes": "2 Yohanes",
    "3 Johannes": "3 Yohanes",
    "Judas": "Yudas",
    "Openbaring": "Wahyu"
}

FIELD_MAP = {
    "Datum": "date",
    "Voorganger": "voorganger",
    "OvD": "ouderling",
    "1e Ontvangst": "eersteO",
    "Muzikanten": "muzikanten",
    "Voorzangers": "voorzangers",
    "Beamer": "beamer",
    "Geluid": "geluid",
    "KND": "knd",
    "Tieners": "tieners",
    "Tikkie link": "qr_link",
    "Titel Prediking": "titelP"
}

IBAN = "NL40.ABNA.0549.3085.12"

NL_MONTHS = [
    "januari","februari","maart","april","mei","juni",
    "juli","augustus","september","oktober","november","december"
]

def normalize_spaces(txt: str) -> str:
    # Quote characters we treat as starting/ending quoted sections
    OPEN_QUOTES = {'“', '‘'}
    CLOSE_QUOTES = {'”', '’'}
    TOGGLE_QUOTES = {'"'}  # straight double quote toggles in/out

    PUNCT = {',', '.', '!', '?', ';', ':'}

    out = []
    in_quote = False

    def last_is_space():
        return len(out) > 0 and out[-1] == ' '

    i = 0
    n = len(txt)

    while i < n:
        ch = txt[i]

        # Track quote state
        if ch in OPEN_QUOTES:
            in_quote = True
            out.append(ch)
            i += 1
            continue
        if ch in CLOSE_QUOTES:
            in_quote = False
            out.append(ch)
            i += 1
            continue
        if ch in TOGGLE_QUOTES:
            in_quote = not in_quote
            out.append(ch)
            i += 1
            continue

        # If we're inside quotes, don't normalize spacing rules—just copy
        if in_quote:
            out.append(ch)
            i += 1
            continue

        # Outside quotes: normalize whitespace to single spaces
        if ch.isspace():
            if not last_is_space() and out:
                out.append(' ')
            i += 1
            continue

        # Outside quotes: punctuation rules
        if ch in PUNCT:
            # Remove space before punctuation
            if last_is_space():
                out.pop()

            out.append(ch)

            # Add a space after punctuation if next char exists and needs it
            j = i + 1
            if j < n:
                nxt = txt[j]
                if (not nxt.isspace()) and (nxt not in TOGGLE_QUOTES) and (nxt not in OPEN_QUOTES) and (nxt not in CLOSE_QUOTES):
                    out.append(' ')
            i += 1
            continue

        # Normal character
        out.append(ch)
        i += 1

    # Final cleanup: strip leading/trailing space
    return ''.join(out).strip()

def is_empty(value):
    # First check for pandas-style NaN / None
    if pd.isna(value):
        return True
    
    # If it's a string, check empty or "nan"
    if isinstance(value, str):
        return value.strip().lower() in ("", "nan")
    
    # If it's a number, check for NaN
    if isinstance(value, (int, float)):
        return math.isnan(value) if isinstance(value, float) else False
    
    return False

def format_date_long_nl(dt: datetime) -> str:
    return f"{dt.day} {NL_MONTHS[dt.month - 1]} {dt.year}"

def format_date_short_nl(dt: datetime) -> str:
    return f"{dt.day}-{dt.month}-{dt.year}"

def set_cell_borders(cell, **borders):
    """
    Set specific cell borders. Use 'nil' with size 0 to remove a border.
    Example: set_cell_borders(cell, bottom=("double", 12))
    """
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    tcBorders = tcPr.first_child_found_in("w:tcBorders")
    if tcBorders is None:
        tcBorders = OxmlElement('w:tcBorders')
        tcPr.append(tcBorders)

    for edge in ("top","left","bottom","right","insideH","insideV"):
        if edge in borders and borders[edge] is not None:
            val, sz = borders[edge]
            el = OxmlElement(f"w:{edge}")
            el.set(qn("w:val"), val)      # 'single','double','nil',...
            el.set(qn("w:sz"), str(sz))   # eighths of a point
            el.set(qn("w:space"), "0")
            el.set(qn("w:color"), "000000")
            existing = tcBorders.find(qn(f"w:{edge}"))
            if existing is not None:
                tcBorders.remove(existing)
            tcBorders.append(el)

def clear_cell_borders(cell):
    set_cell_borders(cell, top=("nil",0), bottom=("nil",0), left=("nil",0), right=("nil",0))

def set_cell_width(cell, width_in_inches: float):
    cell.width = Inches(width_in_inches)

dir_path = os.path.dirname(os.path.realpath(__file__))

NS = {
    "s":   "http://schemas.openxmlformats.org/spreadsheetml/2006/main",
    "r":   "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
    "xdr": "http://schemas.openxmlformats.org/drawingml/2006/spreadsheetDrawing",
    "a":   "http://schemas.openxmlformats.org/drawingml/2006/main",
}




# -------------------------------
# READ EXCEL (Table 1 as vertical list)
# -------------------------------


excelF = pd.read_excel(os.path.join(dir_path + '/file mingguan/', EXCEL_NAME), sheet_name=SHEET_NAME)

table1 = excelF.iloc[1:13, 0:2].copy()  # labels | values
table1 = table1.map(lambda x: x.strip() if isinstance(x, str) else x)

if table1.shape[1] < 2:
    raise ValueError("Table 1 should be two columns (label | value).")

labels = table1.iloc[:, 0].astype(str)
values = table1.iloc[:, 1]

data_map = {}
for lbl, val in zip(labels, values):
    lbl_clean = (lbl or "").strip()
    if lbl_clean in FIELD_MAP:
        data_map[FIELD_MAP[lbl_clean]] = val

required = ["date", "voorganger", "ouderling"]
missing = [k for k in required if k not in data_map]
if missing:
    raise ValueError(f"Missing required fields in Table 1: {missing}")

date_val = data_map["date"]
if isinstance(date_val, datetime):
    dienst_date = date_val
else:
    dienst_date = pd.to_datetime(str(date_val), dayfirst=True).to_pydatetime()

voorganger = str(data_map["voorganger"]).strip()
ouderling = str(data_map["ouderling"]).strip()
eersteO = str(data_map["eersteO"]).strip()
muzikanten = str(data_map["muzikanten"]).strip()
voorzangers = str(data_map["voorzangers"]).strip()
geluid = str(data_map["geluid"]).strip()
beamer = str(data_map["beamer"]).strip()
knd = str(data_map["knd"]).strip()
tieners = str(data_map["tieners"]).strip()
qr_link = str(data_map.get("qr_link", "")).strip()

next_week = dienst_date + timedelta(days=7)
long_date  = format_date_long_nl(dienst_date)
long_datenw  = format_date_long_nl(next_week)
short_date = format_date_short_nl(dienst_date)
titelPred = str(data_map.get("titelP", "")).strip()

def _flatten_text(node):
    if node is None:
        return ""
    if isinstance(node, str):
        return node
    if isinstance(node, list):
        return "".join(_flatten_text(x) for x in node)
    if isinstance(node, dict):
        return _flatten_text(node.get("content"))
    return ""

def _collect_verses_from_chapter(chapter_content, v_from, v_to):
    verses = {v: "" for v in range(v_from, v_to + 1)}
    for paragraph in chapter_content:
        for item in paragraph.get("content", []):
            if item.get("type") == "verse-text":
                verse_id = item.get("verseId")  # e.g., "GEN.1.18"
                try:
                    verse_num = int(verse_id.split(".")[-1])
                except Exception:
                    continue
                if v_from <= verse_num <= v_to:
                    verses[verse_num] += _flatten_text(item.get("content"))
    for v in list(verses.keys()):
        verses[v] = " ".join(verses[v].split())
    return verses

def _to_int_or_raise(x, name):
    try:
        return int(x)
    except Exception:
        raise ValueError(f"{name} moet een geheel getal zijn (gekregen: {x!r})")

def add_verses_to_doc(doc, dutch_book_name, chapter_number, verse_from, verse_to=None, font_name="Calibri", font_size_pt=10):
    """
    - Maps Dutch book name -> abbreviation
    - Opens JSON file named '{ABBR}.{chapter}'
    - Gets verses in inclusive range [verse_from, verse_to]
      * If verse_to is None/""/0 -> uses single verse (verse_to = verse_from)
    - Appends ONE paragraph to 'doc' with superscript verse numbers and normal text (Calibri 10)
    """
    # Normalize inputs (accept strings like "1" or "")
    chapter_number = _to_int_or_raise(chapter_number, "hoofdstuk")
    verse_from = _to_int_or_raise(verse_from, "vers_van")
    if verse_to in (None, "", 0, "0"):
        verse_to = verse_from
    else:
        verse_to = _to_int_or_raise(verse_to, "vers_tot")

    if verse_to < verse_from:
        # If user passed reversed range, swap to be forgiving
        verse_from, verse_to = verse_to, verse_from

    # 1 Map
    try:
        abbr = bible_books_nl[dutch_book_name]
    except KeyError:
        raise ValueError(f"Onbekende boeknaam: {dutch_book_name!r}")

    # 2 Open JSON
    json_path = os.path.join(dir_path, "bible", f"{abbr}.{chapter_number}")
    print(json_path)

    if not os.path.exists(json_path):
        # Try downloading from Dropbox on demand (web deployment)
        _fetch = globals().get('_ensure_bible_file')
        if _fetch:
            try:
                _fetch(f"{abbr}.{chapter_number}")
            except Exception:
                pass
    if not os.path.exists(json_path):
        raise FileNotFoundError(f"Bestand niet gevonden: {json_path}")

    with open(json_path, "r", encoding="utf-8") as f:
        payload = json.load(f)

    try:
        chapter = payload["data"]["chapter"]
        chapter_content = chapter["content"]
    except Exception as e:
        raise ValueError(f"Onverwachte JSON-structuur in {json_path}: {e}")

    # 3 Collect
    verses = _collect_verses_from_chapter(chapter_content, verse_from, verse_to)

    # 4 Write paragraph
    p = doc.add_paragraph()
    for v in range(verse_from, verse_to + 1):
        text = verses.get(v, "")
        if not text:
            continue
        r_num = p.add_run(str(v))
        r_num.font.superscript = True
        r_num.font.name = font_name
        r_num.font.size = Pt(font_size_pt)

        r_txt = p.add_run(" " + text + " ")
        r_txt.font.superscript = False
        r_txt.font.name = font_name
        r_txt.font.size = Pt(font_size_pt)

    return p

# -------------------------------
# CREATE WORD DOC
# -------------------------------

doc = Document()

# --- Set margins FIRST so header table can use full width ---
section = doc.sections[0]
section.top_margin = Cm(0.5)
section.bottom_margin = Cm(1.2)
section.left_margin = Cm(1.5)
section.right_margin = Cm(1.5)

# Set a clean Normal style baseline
doc.styles["Normal"].font.name = "Times New Roman"
doc.styles["Normal"].font.size = Pt(11)

def add_para(text, bold=False, size=11, align=None, keep_next=False, space_before=0, space_after=3, line_spacing=1.0):
    p = doc.add_paragraph()
    if align is not None:
        p.alignment = align
    run = p.add_run(text)
    run.bold = bool(bold)
    run.font.size = Pt(size)
    pf = p.paragraph_format
    pf.space_before = Pt(space_before)
    pf.space_after = Pt(space_after)
    pf.line_spacing = line_spacing
    if keep_next:
        pf.keep_with_next = True
    return p

# === HEADER TABLE (3 columns, 2 rows) ===
header = doc.add_table(rows=2, cols=3)
header.autofit = False  # ensure widths stick

# Remove ALL borders first
for row in header.rows:
    for cell in row.cells:
        clear_cell_borders(cell)

# Stretch to full usable page width so long title fits on one line
usable_width = section.page_width - section.left_margin - section.right_margin
usable_w_in = usable_width / Inches(1)   # float inches

logo_w_in = 1.35
gap_w_in  = 0.08
text_w_in = max(4.5, usable_w_in - logo_w_in - gap_w_in)  # guard minimum

# Set column widths (both rows)
for r in (0, 1):
    set_cell_width(header.cell(r, 0), logo_w_in)
    set_cell_width(header.cell(r, 1), gap_w_in)
    set_cell_width(header.cell(r, 2), text_w_in)

# Merge first column rows and place logo
logo_cell = header.cell(0, 0).merge(header.cell(1, 0))
logo_par = logo_cell.paragraphs[0]
logo_par.alignment = WD_PARAGRAPH_ALIGNMENT.CENTER
logo_cell.vertical_alignment = WD_ALIGN_VERTICAL.CENTER
logo_run = logo_par.add_run()
logo_path = os.path.join(dir_path, LOGO_FILE)
if os.path.exists(logo_path):
    logo_run.add_picture(logo_path, width=Inches(min(1.2, logo_w_in - 0.1)))
else:
    logo_par.add_run("[LOGO ontbreekt: logo.png]").bold = True

# Column 2 (gap) stays empty and borderless
header.cell(0, 1).text = ""
header.cell(1, 1).text = ""

# Column 3, Row 1: titles, tight spacing; ONLY bottom double border
title_cell_top = header.cell(0, 2); title_cell_top.text = ""

p1 = title_cell_top.paragraphs[0]
p1.alignment = WD_PARAGRAPH_ALIGNMENT.LEFT
p1.paragraph_format.space_after = Pt(0)
p1.paragraph_format.line_spacing = 1.0
r1 = p1.add_run("GEREJA KRISTEN INDONESIA NEDERLAND")
r1.bold = True; r1.font.size = Pt(16); r1.font.name = "Bookman Old Style"

p2 = title_cell_top.add_paragraph()
p2.alignment = WD_PARAGRAPH_ALIGNMENT.LEFT
p2.paragraph_format.space_after = Pt(2)
p2.paragraph_format.line_spacing = 1.0
r2 = p2.add_run("Indonesisch Nederlands Christelijke Kerk")
r2.bold = True; r2.italic = True; r2.font.size = Pt(11); r2.font.name = "Times New Roman"

# ONLY border on this cell
set_cell_borders(title_cell_top, bottom=("double",12))

# Column 3, Row 2: remove default blank paragraph, then nested 3-col table (borderless)
title_cell_bot = header.cell(1, 2)

# Remove auto empty paragraph so there's NO blank line
while title_cell_bot.paragraphs:
    p = title_cell_bot.paragraphs[0]._element
    p.getparent().remove(p)

nested = title_cell_bot.add_table(rows=3, cols=3)
nested.autofit = False
for r in nested.rows:
    for c in r.cells:
        clear_cell_borders(c)

# widths inside nested: label 2.25", ":" 0.25", value = remainder
label_w, colon_w = 2.25, 0.25
value_w = max(1.0, text_w_in - label_w - colon_w)
for i in range(3):
    set_cell_width(nested.cell(i,0), label_w)
    set_cell_width(nested.cell(i,1), colon_w)
    set_cell_width(nested.cell(i,2), value_w)

labels = ["Eredienst van zondag", "Voorganger", "Ouderling van Dienst"]
values = [long_date, voorganger, ouderling]

for i in range(3):
    lp = nested.cell(i,0).paragraphs[0]
    lp.paragraph_format.space_after = Pt(0)
    lp.paragraph_format.line_spacing = 1.0
    lrun = lp.add_run(labels[i])
    lrun.font.size = Pt(12); lrun.font.name = "Times New Roman"; lrun.bold = True

    cp = nested.cell(i,1).paragraphs[0]
    cp.alignment = WD_PARAGRAPH_ALIGNMENT.CENTER
    cp.paragraph_format.space_after = Pt(0)
    cp.paragraph_format.line_spacing = 1.0
    crun = cp.add_run(":")
    crun.font.size = Pt(12); crun.font.name = "Times New Roman"; crun.bold = True

    vp = nested.cell(i,2).paragraphs[0]
    vp.paragraph_format.space_after = Pt(0)
    vp.paragraph_format.line_spacing = 1.0
    vrun = vp.add_run(values[i])
    vrun.font.size = Pt(12); vrun.font.name = "Times New Roman"; vrun.bold = True

# -------------------------------
# COLLECTION BLOCK (ADJUSTED ONLY)
# -------------------------------

# ~ small space before the block (tight)
sp = doc.add_paragraph()
sp.paragraph_format.space_before = Pt(6)
sp.paragraph_format.space_after = Pt(0)
sp.paragraph_format.line_spacing = 1.0

# Two-column table: left=text, right=QR; only top/bottom borders
coll_tbl = doc.add_table(rows=1, cols=2)
coll_tbl.autofit = False

# Full usable width (match header width)
left_w_in = max(1.0, usable_w_in - 2.0)  # right column fixed ~2.0" for QR
right_w_in = 2.0
set_cell_width(coll_tbl.cell(0,0), left_w_in)
set_cell_width(coll_tbl.cell(0,1), right_w_in)

# Clear borders first
for cell in coll_tbl.rows[0].cells:
    clear_cell_borders(cell)
# Apply only top & bottom borders to both cells (continuous line across the row)
for cell in coll_tbl.rows[0].cells:
    set_cell_borders(cell, top=("single",12), bottom=("single",12))

# LEFT CELL: Calibri 9 (title + intro + 3 numbered lines)
left_cell = coll_tbl.cell(0,0)
while left_cell.paragraphs:
    pe = left_cell.paragraphs[0]._element
    pe.getparent().remove(pe)

def add_calibri9_par(cell, text, bold=False):
    p = cell.add_paragraph()
    pf = p.paragraph_format
    pf.space_before = Pt(0)
    pf.space_after = Pt(0)
    pf.line_spacing = 1.0
    run = p.add_run(text)
    run.font.name = "Calibri"
    run.font.size = Pt(9)
    run.bold = bold
    return p

add_calibri9_par(left_cell, f"Regionale collecte zondag {short_date} ({REGIO})", bold=True)
add_calibri9_par(left_cell, "Tevens heeft u de gelegenheid om uw dankoffer tijdens deze fysieke eredienst te geven op drie manieren:")

items = [
    "Deponeren in de collectezak;",
    f"Overmaken naar rekeningnummer van GKIN Amstelveen,\n    IBAN: {IBAN}, o.v.v. “Collecte {short_date}”",
    f"Gebruik te maken van de QR-code, of betaalverzoek link: {qr_link}",
]
for idx, txt in enumerate(items, start=1):
    add_calibri9_par(left_cell, f"{idx}. {txt}")

# RIGHT CELL: Calibri 9 QR placeholder only
right_cell = coll_tbl.cell(0,1)
while right_cell.paragraphs:
    pe = right_cell.paragraphs[0]._element
    pe.getparent().remove(pe)
p_qr = right_cell.add_paragraph()
p_qr.alignment = WD_PARAGRAPH_ALIGNMENT.CENTER
rq = p_qr.add_run("\n\n\nQR CODE\nPLACEHOLDER\n")
rq.font.name = "Calibri"
rq.font.size = Pt(9)
rq.bold = True

# -------------------------------
# TABLE 2 & 3 parsing (once)
# -------------------------------

raw_table2 = excelF.iloc[14:20, 1:9]
table2 = raw_table2[1:].copy()
table2.columns = raw_table2.iloc[0]
table2 = table2.reset_index(drop=True)

raw_table3 = excelF.iloc[21:29, 1:17]
table3 = raw_table3[1:].copy()
table3.columns = raw_table3.iloc[0]
table3 = table3.reset_index(drop=True)

# ---------- Parse Table 3 into songs with ordered lyric cells ----------
import re
from math import ceil

def _norm(name): return str(name).strip().lower()
cols = list(table3.columns)

# Required columns (case-insensitive)
col_map = { _norm(c): c for c in cols }
for needed in ["boek","nummer","versen","titel"]:
    if needed not in col_map:
        raise ValueError(f'Missing column: "{needed}" in table 3')
boek_col   = col_map["boek"]
nummer_col = col_map["nummer"]
versen_col = col_map["versen"]
titel_col  = col_map["titel"]

# YouTube column = the column immediately after "Titel"
titel_idx   = cols.index(titel_col)
youtube_col = cols[titel_idx + 1] if titel_idx + 1 < len(cols) else None

# Lyric columns: only those named like "Cell 1..10" (avoid pulling YouTube)
cell_pat = re.compile(r'^\s*cell\s*(\d+)\s*$', re.I)
def _cell_index(name):
    m = cell_pat.match(str(name).strip())
    return int(m.group(1)) if m else 9999

lyric_cols = [c for c in cols if cell_pat.match(str(c).strip())]
lyric_cols = sorted(lyric_cols, key=_cell_index)

songs = []
for _, row in table3.iterrows():
    yt = ""
    if youtube_col is not None:
        v = row.get(youtube_col, "")
        yt = ("" if pd.isna(v) else str(v)).strip()

    cells = []
    for c in lyric_cols:
        v = row.get(c, None)
        if isinstance(v, str) and v.strip():
            cells.append(v.strip())

    songs.append({
        "Boek":    str(row.get(boek_col, "")).strip(),
        "Nummer":  str(row.get(nummer_col, "")).strip(),
        "Versen":  str(row.get(versen_col, "")).strip(),
        "Titel":   str(row.get(titel_col, "")).strip(),
        "youtube": yt,
        "cells":   cells,
    })

# ---------- Helpers ----------
def add_calibri_paragraph(container, text, size_pt=10, bold=False, keep_next=False):
    """Add a Calibri paragraph with 0/0 spacing; keep_next is opt-in."""
    p = container.add_paragraph()
    pf = p.paragraph_format
    pf.space_before = Pt(0)
    pf.space_after  = Pt(0)
    pf.line_spacing = 1.0
    pf.keep_with_next = bool(keep_next)
    pf.keep_together  = True
    r = p.add_run(text)
    r.font.name = "Calibri"
    r.font.size = Pt(size_pt)
    r.bold = bold
    return p

_num_re = re.compile(r'(?m)^\s*\d+([.)]|\s)\s*|(?<!\w)\d{1,3}(?!\w)')

def _cell_has_number(text: str) -> bool:
    """True if any line in the cell starts with a number OR a standalone number appears."""
    return bool(_num_re.search(text or ""))

def _clear_cell(cell):
    while cell.paragraphs:
        pe = cell.paragraphs[0]._element
        pe.getparent().remove(pe)

def _set_row_cant_split(row):
    """Prevent a row from breaking across pages."""
    tr = row._tr
    trPr = tr.get_or_add_trPr()
    cant = OxmlElement('w:cantSplit')
    trPr.append(cant)

def _write_cell(cell, text: str):
    """Write text into a table cell (Calibri 10), no extra blank lines added."""
    _clear_cell(cell)
    add_calibri_paragraph(cell, text, size_pt=10, bold=False, keep_next=False)

# ---------- Main callable ----------
def add_song_to_doc(song_no: int, intro_text=""):
    """Render the specified song (1–7) into the document at the current position."""
    if not (1 <= song_no <= len(songs)):
        raise ValueError(f"Song number out of range: {song_no} (have {len(songs)} songs)")
    song = songs[song_no - 1]

    if song["Boek"] == "OPW":
        song["Boek"] = "Opwekking"

    # Title: [optional intro] Boek Nummer: Versen "Titel" — bold Calibri 10
    prefix = (intro_text.strip() + " ") if intro_text.strip() else ""
    if song["Boek"] == "Kosong":
        title_text = f'{prefix} "{song["Titel"]}"'
    else:
        if is_empty(song["Nummer"]):
            title_text = f'{prefix}{song["Boek"]} "{song["Titel"]}"'
        else:
            if is_empty(song["Versen"]):
                title_text = f'{prefix}{song["Boek"]} {song["Nummer"]} "{song["Titel"]}"'
            else:
                title_text = f'{prefix}{song["Boek"]} {song["Nummer"]}: {song["Versen"]} "{song["Titel"]}"'
    add_calibri_paragraph(doc, title_text, size_pt=10, bold=True, keep_next=True)  # keep title with its table

    n = len(song["cells"])
    if n == 0:
        return

    # Helper to safely append a single line break if not already there
    def _append_final_linebreak(text: str) -> str:
        if text is None:
            return ""
        text = str(text)
        return text if text.endswith("\n") else text + "\n"

    if n <= 2:
        tbl = doc.add_table(rows=1, cols=n)
        tbl.autofit = True
        for row in tbl.rows:
            _set_row_cant_split(row)
            for cell in row.cells:
                clear_cell_borders(cell)
        for j, text in enumerate(song["cells"]):
            _write_cell(tbl.cell(0, j), _append_final_linebreak(text))
    else:
        left_len = ceil(n / 2)
        right_len = n - left_len
        rows_needed = left_len

        tbl = doc.add_table(rows=rows_needed, cols=2)
        tbl.autofit = True
        for r_idx, row in enumerate(tbl.rows):
            _set_row_cant_split(row)
            for cell in row.cells:
                clear_cell_borders(cell)

        for i in range(left_len):
            _write_cell(tbl.cell(i, 0), _append_final_linebreak(song["cells"][i]))
        for i in range(right_len):
            _write_cell(tbl.cell(i, 1), _append_final_linebreak(song["cells"][left_len + i]))

    add_calibri_paragraph(doc, "", size_pt=10, bold=False, keep_next=False)

# ------- Content after collection block -------
add_calibri_paragraph(doc, "", size_pt=10, bold=False, keep_next=False)
add_calibri_paragraph(doc, "", size_pt=10, bold=False, keep_next=False)

# Intochtslied (song 1) with label on same line
add_song_to_doc(1, "Intochtslied:")

# 1. Stil gebed (de gemeente staat) — inline bold + italic, Calibri 10
p = doc.add_paragraph()
pf = p.paragraph_format
pf.space_before = Pt(6); pf.space_after = Pt(6); pf.line_spacing = 1.0
r1 = p.add_run("1. Stil gebed ")
r1.font.name = "Calibri"; r1.font.size = Pt(10); r1.bold = True
r2 = p.add_run("( de gemeente staat )")
r2.font.name = "Calibri"; r2.font.size = Pt(10); r2.bold = False; r2.italic = True

from docx.shared import Cm, Pt

# --- Element 2: Votum en groet (staande) ---

# Title line: bold + bracket note italic, Calibri 10, 0/0 spacing
p = doc.add_paragraph()
pf = p.paragraph_format
pf.space_before = Pt(0); pf.space_after = Pt(0); pf.line_spacing = 1.0
run = p.add_run("2. Votum en groet ")
run.font.name = "Calibri"; run.font.size = Pt(10); run.bold = True
run = p.add_run("(staande)")
run.font.name = "Calibri"; run.font.size = Pt(10); run.bold = False; run.italic = True
pf.keep_with_next = True  # keep title with the following lines

# Common tab stop so texts after labels align in a clean column
tab_pos = Cm(1.5)  # adjust if you want more/less indent

# Helper to add a role line (Voorganger/Gemeente) with aligned tab
# (still inline code — just repeat these few lines per row you need)
# --- Voorganger line 1 ---
p = doc.add_paragraph()
pf = p.paragraph_format
pf.space_before = Pt(0); pf.space_after = Pt(0); pf.line_spacing = 1.0
pf.keep_with_next = True
pf.tab_stops.add_tab_stop(tab_pos)
r = p.add_run("\tVoorganger: ")
r.font.name = "Calibri"; r.font.size = Pt(10)
r = p.add_run("Onze hulp is in de Naam van de HEER,")
r.font.name = "Calibri"; r.font.size = Pt(10)

# --- Gemeente line 1 ---
p = doc.add_paragraph()
pf = p.paragraph_format
pf.space_before = Pt(0); pf.space_after = Pt(0); pf.line_spacing = 1.0
pf.keep_with_next = True
pf.tab_stops.add_tab_stop(tab_pos)
r = p.add_run("\t\tGemeente: ")
r.font.name = "Calibri"; r.font.size = Pt(10)
r = p.add_run("die hemel en aarde gemaakt heeft.")
r.font.name = "Calibri"; r.font.size = Pt(10)

# (Add more lines as needed, reusing the same pattern)
# --- Voorganger line 2 ---
p = doc.add_paragraph()
pf = p.paragraph_format
pf.space_before = Pt(0); pf.space_after = Pt(0); pf.line_spacing = 1.0
pf.keep_with_next = True
pf.tab_stops.add_tab_stop(tab_pos)
r = p.add_run("\tVoorganger: ")
r.font.name = "Calibri"; r.font.size = Pt(10)
r = p.add_run("die trouw houdt tot in eeuwigheid")
r.font.name = "Calibri"; r.font.size = Pt(10)

# --- Gemeente line 2 ---
p = doc.add_paragraph()
pf = p.paragraph_format
pf.space_before = Pt(0); pf.space_after = Pt(0); pf.line_spacing = 1.0
pf.keep_with_next = True
pf.tab_stops.add_tab_stop(tab_pos)
r = p.add_run("\t\tGemeente: ")
r.font.name = "Calibri"; r.font.size = Pt(10)
r = p.add_run("en niet laat varen het werk van Zijn handen.")
r.font.name = "Calibri"; r.font.size = Pt(10)

# (Add more lines as needed, reusing the same pattern)
# --- Voorganger line 3 ---
p = doc.add_paragraph()
pf = p.paragraph_format
pf.space_before = Pt(0); pf.space_after = Pt(0); pf.line_spacing = 1.0
pf.keep_with_next = True
pf.tab_stops.add_tab_stop(tab_pos)
r = p.add_run("\tVoorganger: ")
r.font.name = "Calibri"; r.font.size = Pt(10)
r = p.add_run("Genade zij u en vrede van God, onze Vader, en van de Heer Jezus Christus, in de gemeenschap van de Heilige Geest.")
r.font.name = "Calibri"; r.font.size = Pt(10)

# --- Gemeente line 3 ---
p = doc.add_paragraph()
pf = p.paragraph_format
pf.space_before = Pt(0); pf.space_after = Pt(6); pf.line_spacing = 1.0
pf.keep_with_next = False
pf.tab_stops.add_tab_stop(tab_pos)
r = p.add_run("\t\tGemeente:")
r.font.name = "Calibri"; r.font.size = Pt(10)
r = p.add_run("Amen.")
r.font.name = "Calibri"; r.font.size = Pt(10)

# --- Element 3: Aanvangstekst (zittende) from table2 row 1 ---
def _get(table, row_idx, names):
    for name in names:
        if name in table.columns:
            val = table.loc[row_idx, name]
            return ("" if pd.isna(val) else str(val)).strip()
    return ""

boek = _get(table2, 0, ["Boek", "boek"])
hs   = _get(table2, 0, ["H.s.", "H.s", "HS", "H.S.", "H.S"])
vfrom = _get(table2, 0, ["Vers van", "Vers Van", "vers van", "Van"])
vto   = _get(table2, 0, ["Vers tot", "Vers Tot", "vers tot", "Tot"])
boek2 = _get(table2, 0, ["Boek 2", "boek 2"])
hs2   = _get(table2, 0, ["H.s. 2", "H.s 2", "HS 2", "H.S. 2", "H.S 2"])
vfrom2 = _get(table2, 0, ["Vers van 2", "Vers Van 2", "vers van 2", "Van 2"])
vto2   = _get(table2, 0, ["Vers tot 2", "Vers Tot 2", "vers tot 2", "Tot 2"])

if is_empty(vfrom):
    elem3_text = f"3. Aanvangstekst (zittende)   {boek} {hs}"
elif is_empty(vto):    
    elem3_text = f"3. Aanvangstekst (zittende)   {boek} {hs}:{vfrom}"
else:
    elem3_text = f"3. Aanvangstekst (zittende)   {boek} {hs}:{vfrom}-{vto}."

if is_empty(boek2) == False:
    if vfrom2 and str(vfrom2).strip():
        if vto2 and str(vto2).strip():
            elem3_text = elem3_text + f" en {boek2} {hs2}:{vfrom2}-{vto2}"
        else:
            elem3_text = elem3_text + f" en {boek2} {hs2}:{vfrom2}"
    else:
        elem3_text = elem3_text + f" en {boek2}"

p = doc.add_paragraph()
pf = p.paragraph_format
pf.space_before = Pt(0); pf.space_after = Pt(6); pf.line_spacing = 1.0
run = p.add_run(elem3_text); run.font.name = "Calibri"; run.font.size = Pt(10); run.bold = True

# --- Element 4: 2e lied ---
add_song_to_doc(2, "4.")

# --- Element 5: De kinderen ---
p = doc.add_paragraph()
pf = p.paragraph_format
pf.space_before = Pt(0); pf.space_after = Pt(6); pf.line_spacing = 1.0; pf.space_before = Pt(6)
if is_empty(tieners):
    r = p.add_run("5. De kinderen gaan naar de kindernevendienst. Er is geen tienerdienst.")
else:
    r = p.add_run("5. De kinderen gaan naar de kindernevendienst. Tieners gaan naar de tienerdienst.")
r.font.name = "Calibri"; r.font.size = Pt(10); r.bold=True

# --- Element 6: Gebed van toenadering ---
p = doc.add_paragraph()
pf = p.paragraph_format
pf.space_before = Pt(0); pf.space_after = Pt(6); pf.line_spacing = 1.0
r = p.add_run("6. Gebed van toenadering.")
r.font.name = "Calibri"; r.font.size = Pt(10); r.bold=True

# --- Element 7: 3e lied ---
add_song_to_doc(3, "7.")

# --- Element 8: Genadeverkondiging from table2 row 2 ---
boek = _get(table2, 1, ["Boek", "boek"])
hs   = _get(table2, 1, ["H.s.", "H.s", "HS", "H.S.", "H.S"])
vfrom = _get(table2, 1, ["Vers van", "Vers Van", "vers van", "Van"])
vto   = _get(table2, 1, ["Vers tot", "Vers Tot", "vers tot", "Tot"])
boek2 = _get(table2, 1, ["Boek 2", "boek 2"])
hs2   = _get(table2, 1, ["H.s. 2", "H.s 2", "HS 2", "H.S. 2", "H.S 2"])
vfrom2 = _get(table2, 1, ["Vers van 2", "Vers Van 2", "vers van 2", "Van 2"])
vto2   = _get(table2, 1, ["Vers tot 2", "Vers Tot 2", "vers tot 2", "Tot 2"])

if is_empty(vfrom):
    elem8_text = f"8. Genadeverkondiging   {boek} {hs}"
elif is_empty(vto):    
    elem8_text = f"8. Genadeverkondiging   {boek} {hs}:{vfrom}"
else:
    elem8_text = f"8. Genadeverkondiging   {boek} {hs}:{vfrom}-{vto}."

if is_empty(boek2) == False:
    if vfrom2 and str(vfrom2).strip():
        if vto2 and str(vto2).strip():
            elem8_text = elem8_text + f" en {boek2} {hs2}:{vfrom2}-{vto2}"
        else:
            elem8_text = elem8_text + f" en {boek2} {hs2}:{vfrom2}"
    else:
        elem8_text = elem8_text + f" en {boek2}"

p = doc.add_paragraph()
pf = p.paragraph_format
pf.space_before = Pt(0); pf.space_after = Pt(6); pf.line_spacing = 1.0; pf.keep_with_next=False
run = p.add_run(elem8_text); run.font.name = "Calibri"; run.font.size = Pt(10); run.bold = True

# --- Element 9: Gods Woord voor levensvernieuwing from table2 row 3 ---
boek = _get(table2, 2, ["Boek", "boek"])
hs   = _get(table2, 2, ["H.s.", "H.s", "HS", "H.S.", "H.S"])
vfrom = _get(table2, 2, ["Vers van", "Vers Van", "vers van", "Van"])
vto   = _get(table2, 2, ["Vers tot", "Vers Tot", "vers tot", "Tot"])
boek2 = _get(table2, 2, ["Boek 2", "boek 2"])
hs2   = _get(table2, 2, ["H.s. 2", "H.s 2", "HS 2", "H.S. 2", "H.S 2"])
vfrom2 = _get(table2, 2, ["Vers van 2", "Vers Van 2", "vers van 2", "Van 2"])
vto2   = _get(table2, 2, ["Vers tot 2", "Vers Tot 2", "vers tot 2", "Tot 2"])

if is_empty(vfrom):
    elem9_text = f"9. Gods Woord voor levensvernieuwing   {boek} {hs}"
elif is_empty(vto):    
    elem9_text = f"9. Gods Woord voor levensvernieuwing   {boek} {hs}:{vfrom}"
else:
    elem9_text = f"9. Gods Woord voor levensvernieuwing   {boek} {hs}:{vfrom}-{vto}"

if is_empty(boek2) == False:
    if vfrom2 and str(vfrom2).strip():
        if vto2 and str(vto2).strip():
            elem9_text = elem9_text + f" en {boek2} {hs2}:{vfrom2}-{vto2}"
        else:
            elem9_text = elem9_text + f" en {boek2} {hs2}:{vfrom2}"
    else:
        elem9_text = elem9_text + f" en {boek2}"


p = doc.add_paragraph()
pf = p.paragraph_format
pf.space_before = Pt(0); pf.space_after = Pt(6); pf.line_spacing = 1.0; pf.keep_with_next=False
run = p.add_run(elem9_text); run.font.name = "Calibri"; run.font.size = Pt(10); run.bold = True

# --- Element 10: 4e lied ---
add_song_to_doc(4, "10. ")

# --- Element 11: Gebed om verlichting met de Heilige Geest ---
p = doc.add_paragraph()
pf = p.paragraph_format
pf.space_before = Pt(0); pf.space_after = Pt(6); pf.line_spacing = 1.0; pf.space_before = Pt(6); pf.keep_with_next=False
r = p.add_run("11. Gebed om verlichting met de Heilige Geest")
r.font.name = "Calibri"; r.font.size = Pt(10); r.bold=True

# --- Element 12: Schriftlezing ---
boek = _get(table2, 3, ["Boek", "boek"])
hs   = _get(table2, 3, ["H.s.", "H.s", "HS", "H.S.", "H.S"])
vfrom = _get(table2, 3, ["Vers van", "Vers Van", "vers van", "Van"])
vto   = _get(table2, 3, ["Vers tot", "Vers Tot", "vers tot", "Tot"])
boek2 = _get(table2, 3, ["Boek 2", "boek 2"])
hs2   = _get(table2, 3, ["H.s. 2", "H.s 2", "HS 2", "H.S. 2", "H.S 2"])
vfrom2 = _get(table2, 3, ["Vers van 2", "Vers Van 2", "vers van 2", "Van 2"])
vto2   = _get(table2, 3, ["Vers tot 2", "Vers Tot 2", "vers tot 2", "Tot 2"])

if vfrom and str(vfrom).strip():
    if vto and str(vto).strip():
        elem12_text = f"12. Schriftlezing   {boek} {hs}:{vfrom}-{vto}"
    else:
        elem12_text = f"12. Schriftlezing   {boek} {hs}:{vfrom}"
else:
    elem12_text = f"12. Schriftlezing   {boek}"

if is_empty(boek2) == False:
    if vfrom2 and str(vfrom2).strip():
        if vto2 and str(vto2).strip():
            elem12_text = elem12_text + f" en {boek2} {hs2}:{vfrom2}-{vto2}"
        else:
            elem12_text = elem12_text + f" en {boek2} {hs2}:{vfrom2}"
    else:
        elem12_text = elem12_text + f" en {boek2}"

p = doc.add_paragraph()
pf = p.paragraph_format
pf.space_before = Pt(6); pf.space_after = Pt(0); pf.line_spacing = 1.0; pf.keep_with_next=True
run = p.add_run(elem12_text); run.font.name = "Calibri"; run.font.size = Pt(10); run.bold = True
run2 = p.add_run(" ( gelezen door 2 gemeenteleden )"); run2.font.name = "Calibri"; run2.font.size = Pt(10); run2.bold = False; r2.italic = True

# --- Voorganger line 1 ---
p = doc.add_paragraph()
pf = p.paragraph_format
pf.space_before = Pt(0); pf.space_after = Pt(0); pf.line_spacing = 1.0
pf.keep_with_next = True
pf.tab_stops.add_tab_stop(tab_pos)
r = p.add_run("\tVoorganger: ")
r.font.name = "Calibri"; r.font.size = Pt(10)
r = p.add_run("Gelukkig zij die luisteren naar het Woord van God en ernaar leven.")
r.font.name = "Calibri"; r.font.size = Pt(10)

# --- Gemeente line 1 ---
p = doc.add_paragraph()
pf = p.paragraph_format
pf.space_before = Pt(0); pf.space_after = Pt(6); pf.line_spacing = 1.0
pf.keep_with_next = False
pf.tab_stops.add_tab_stop(tab_pos)
r = p.add_run("\t\tGemeente zingt: ")
r.font.name = "Calibri"; r.font.size = Pt(10)
r = p.add_run(f"‘{HMH} (3x)’ uit KJ 473b")
r.font.name = "Calibri"; r.font.size = Pt(10)

# --- Element 13: Prediking ---
p = doc.add_paragraph()
pf = p.paragraph_format
pf.space_before = Pt(0); pf.space_after = Pt(6); pf.line_spacing = 1.0; pf.keep_with_next=False
r = p.add_run(f'13. Prediking   "{titelPred}".')
r.font.name = "Calibri"; r.font.size = Pt(10); r.bold=True

# --- Element 14: Meditatief Moment ---
p = doc.add_paragraph()
pf = p.paragraph_format
pf.space_before = Pt(0); pf.space_after = Pt(6); pf.line_spacing = 1.0; pf.keep_with_next=False
r = p.add_run('14. Meditatief moment / Voordracht.')
r.font.name = "Calibri"; r.font.size = Pt(10); r.bold=True

# --- Element 15: APOSTOLISCHE GELOOFSBELIJDENIS  ---
p = doc.add_paragraph()
pf = p.paragraph_format
pf.space_before = Pt(0); pf.space_after = Pt(6); pf.line_spacing = 1.0; pf.keep_with_next=False
r = p.add_run('15. APOSTOLISCHE GELOOFSBELIJDENIS uit te spreken door een ouderling. ')
r.font.name = "Calibri"; r.font.size = Pt(10); r.bold=True
r2 = p.add_run(' (de gemeente gaat staan)')
r2.font.name = "Calibri"; r2.font.size = Pt(10); r2.bold=False; r2.italic=True

# --- Element 16: 5e lied  ---
add_song_to_doc(5, "16. ")

# --- Element 17: Voorbede ---
p = doc.add_paragraph()
pf = p.paragraph_format
pf.space_before = Pt(0); pf.space_after = Pt(6); pf.line_spacing = 1.0; pf.keep_with_next=False
r = p.add_run('17. Voorbede ( eindigend met het Onze Vader )')
r.font.name = "Calibri"; r.font.size = Pt(10); r.bold=True

# --- Element 18: Dankoffer ---
boek = _get(table2, 4, ["Boek", "boek"])
hs   = _get(table2, 4, ["H.s.", "H.s", "HS", "H.S.", "H.S"])
vfrom = _get(table2, 4, ["Vers van", "Vers Van", "vers van", "Van"])
vto   = _get(table2, 4, ["Vers tot", "Vers Tot", "vers tot", "Tot"])

p = doc.add_paragraph()
pf = p.paragraph_format
pf.space_before = Pt(0); pf.space_after = Pt(6); pf.line_spacing = 1.0; pf.keep_with_next=True
r = p.add_run('18. Dankoffer')
r.font.name = "Calibri"; r.font.size = Pt(10); r.bold=True

p = doc.add_paragraph()
pf = p.paragraph_format
pf.space_before = Pt(0); pf.space_after = Pt(0); pf.line_spacing = 1.0
pf.keep_with_next = True
pf.tab_stops.add_tab_stop(tab_pos)
r = p.add_run("\tOuderling: Gemeente, wij krijgen nu de gelegenheid ons dankoffer aan God te brengen.  ")
r.font.name = "Calibri"; r.font.size = Pt(10)

p = doc.add_paragraph()
pf = p.paragraph_format
pf.space_before = Pt(0); pf.space_after = Pt(0); pf.line_spacing = 1.0
pf.keep_with_next = True
pf.tab_stops.add_tab_stop(tab_pos)

if vfrom and str(vfrom).strip():
    if vto and str(vto).strip():
        elem13_text = f"En we gedenken daarbij de woorden uit {boek} {hs}:{vfrom}-{vto}."
    else:
        elem13_text = f"En we gedenken daarbij de woorden uit {boek} {hs}:{vfrom}."
else:
    elem13_text = f"En we gedenken daarbij de woorden uit {boek}."

add_verses_to_doc(doc, boek, hs, vfrom, vto)

r = p.add_run(f"\t\t{elem13_text}")
r.font.name = "Calibri"; r.font.size = Pt(10)

p = doc.add_paragraph()
pf = p.paragraph_format
pf.space_before = Pt(0); pf.space_after = Pt(0); pf.line_spacing = 1.0
pf.keep_with_next = True
pf.tab_stops.add_tab_stop(tab_pos)
r = p.add_run(f"\t\tHet dankoffer kunt u geven d.m.v. één van de drie genoemde manieren.")
r.font.name = "Calibri"; r.font.size = Pt(10)

p = doc.add_paragraph()
pf = p.paragraph_format
pf.space_before = Pt(0); pf.space_after = Pt(0); pf.line_spacing = 1.0
pf.keep_with_next = True
pf.tab_stops.add_tab_stop(tab_pos)
r = p.add_run("")
r.font.name = "Calibri"; r.font.size = Pt(10)

add_song_to_doc(6,"Tijdens het dankoffer zingen wij samen uit ")

# --- Element 19: 7e lied  ---
add_song_to_doc(7, "19. ")

# --- Element 20: Zending en Zegen  ---
p = doc.add_paragraph()
pf = p.paragraph_format
pf.space_before = Pt(0); pf.space_after = Pt(0); pf.line_spacing = 1.0; pf.keep_with_next=True
r = p.add_run('20. Zending en Zegen ')
r.font.name = "Calibri"; r.font.size = Pt(10); r.bold=True
r2 = p.add_run(' (de gemeente blijft staan)')
r2.font.name = "Calibri"; r2.font.size = Pt(10); r2.bold=False; r2.italic=True

# --- Voorganger line 1 ---
p = doc.add_paragraph()
pf = p.paragraph_format
pf.space_before = Pt(0); pf.space_after = Pt(0); pf.line_spacing = 1.0
pf.keep_with_next = True
pf.tab_stops.add_tab_stop(tab_pos)
r = p.add_run("\tVoorganger: ")
r.font.name = "Calibri"; r.font.size = Pt(10)
r = p.add_run("Verheft uw harten tot de Heer.")
r.font.name = "Calibri"; r.font.size = Pt(10)

# --- Gemeente line 1 ---
p = doc.add_paragraph()
pf = p.paragraph_format
pf.space_before = Pt(0); pf.space_after = Pt(0); pf.line_spacing = 1.0
pf.keep_with_next = True
pf.tab_stops.add_tab_stop(tab_pos)
r = p.add_run("\t\tGemeente: ")
r.font.name = "Calibri"; r.font.size = Pt(10)
r = p.add_run("Wij zijn met ons hart bij de Heer")
r.font.name = "Calibri"; r.font.size = Pt(10)

# --- Voorganger line 2 ---
p = doc.add_paragraph()
pf = p.paragraph_format
pf.space_before = Pt(0); pf.space_after = Pt(0); pf.line_spacing = 1.0
pf.keep_with_next = True
pf.tab_stops.add_tab_stop(tab_pos)
r = p.add_run("\tVoorganger: ")
r.font.name = "Calibri"; r.font.size = Pt(10)
r = p.add_run("Wees getuigen van Christus")
r.font.name = "Calibri"; r.font.size = Pt(10)

# --- Gemeente line 2 ---
p = doc.add_paragraph()
pf = p.paragraph_format
pf.space_before = Pt(0); pf.space_after = Pt(0); pf.line_spacing = 1.0
pf.keep_with_next = True
pf.tab_stops.add_tab_stop(tab_pos)
r = p.add_run("\t\tGemeente: ")
r.font.name = "Calibri"; r.font.size = Pt(10)
r = p.add_run("Lof aan God")
r.font.name = "Calibri"; r.font.size = Pt(10)

# --- Voorganger line 3 ---
p = doc.add_paragraph()
pf = p.paragraph_format
pf.space_before = Pt(0); pf.space_after = Pt(0); pf.line_spacing = 1.0
pf.keep_with_next = True
pf.tab_stops.add_tab_stop(tab_pos)
r = p.add_run("\tVoorganger: ")
r.font.name = "Calibri"; r.font.size = Pt(10)
r = p.add_run("Geloofd zij de Heer")
r.font.name = "Calibri"; r.font.size = Pt(10)

# --- Gemeente line 3 ---
p = doc.add_paragraph()
pf = p.paragraph_format
pf.space_before = Pt(0); pf.space_after = Pt(0); pf.line_spacing = 1.0
pf.keep_with_next = True
pf.tab_stops.add_tab_stop(tab_pos)
r = p.add_run("\t\tGemeente: ")
r.font.name = "Calibri"; r.font.size = Pt(10)
r = p.add_run("Nu en voor altijd")
r.font.name = "Calibri"; r.font.size = Pt(10)

# --- Voorganger line 4 ---
p = doc.add_paragraph()
pf = p.paragraph_format
pf.space_before = Pt(0); pf.space_after = Pt(6); pf.line_spacing = 1.0
pf.keep_with_next = False
pf.tab_stops.add_tab_stop(tab_pos)
r = p.add_run("\tVoorganger: ")
r.font.name = "Calibri"; r.font.size = Pt(10)
r = p.add_run("Ontvang nu de zegen van de Heer:……")
r.font.name = "Calibri"; r.font.size = Pt(10)

# --- Element 21: Amen ---
p = doc.add_paragraph()
pf = p.paragraph_format
pf.space_before = Pt(0); pf.space_after = Pt(6); pf.line_spacing = 1.0; pf.keep_with_next=False
r = p.add_run('21. Gemeente zingt : “Amen, amen, amen”.')
r.font.name = "Calibri"; r.font.size = Pt(10); r.bold=True

# --- Element 22: Stil gebed ---
p = doc.add_paragraph()
pf = p.paragraph_format
pf.space_before = Pt(0); pf.space_after = Pt(6); pf.line_spacing = 1.0; pf.keep_with_next=False
r = p.add_run('22. Stil gebed.')
r.font.name = "Calibri"; r.font.size = Pt(10); r.bold=True

# -------------------------------
# Save
# -------------------------------
out_name = f"LiturgieA {short_date}_{REGIO}.docx"
out_path = os.path.join(dir_path + '/file mingguan/', out_name)
doc.save(out_path)
print(f"Word document generated: {out_path}")

# -------------------------------
# CREATE WORD DOC2 for Music Team
# -------------------------------

from docx.shared import Pt, RGBColor
from docx.oxml.ns import qn
from openpyxl import load_workbook
from openpyxl.worksheet.table import Table
from typing import List, Tuple, Dict

excel_path = os.path.join(dir_path, EXCEL_NAME)

from docx.shared import Inches

doc2 = Document()

def find_sheet_case_insensitive(wb, target: str):
    for name in wb.sheetnames:
        if name.strip().lower() == target.strip().lower():
            return wb[name]
    # fallback: closest match by containment
    for name in wb.sheetnames:
        if target.strip().lower() in name.strip().lower():
            return wb[name]
    return None

def table_bounds(ws, table: Table) -> Tuple[int, int, int, int]:
    """returns (min_row, min_col, max_row, max_col) 1-based indices"""
    ref = table.ref  # e.g., "B3:E15"
    from openpyxl.utils import range_boundaries
    min_col, min_row, max_col, max_row = range_boundaries(ref)
    return (min_row, min_col, max_row, max_col)

def read_table_to_df(ws, table: Table) -> pd.DataFrame:
    (min_row, min_col, max_row, max_col) = table_bounds(ws, table)
    data = []
    for r in ws.iter_rows(min_row=min_row, max_row=max_row, min_col=min_col, max_col=max_col, values_only=True):
        data.append(list(r))
    if not data:
        return pd.DataFrame()
    header = data[0]
    rows = data[1:]
    df = pd.DataFrame(rows, columns=header)
    # drop completely empty rows
    df = df.dropna(how='all')
    return df

def detect_blocks(ws, min_non_empty_per_row:int=1) -> List[Tuple[int,int,int,int]]:
    """
    Heuristically detect rectangular data blocks ("tables") separated by blank rows.
    Returns a list of (min_row, min_col, max_row, max_col), 1-based.
    """
    max_row = ws.max_row
    max_col = ws.max_column

    def row_non_empty_count(r):
        cnt = 0
        for c in range(1, max_col+1):
            v = ws.cell(r, c).value
            if v is not None and str(v).strip() != "":
                cnt += 1
        return cnt

    blocks = []
    r = 1
    while r <= max_row:
        while r <= max_row and row_non_empty_count(r) < min_non_empty_per_row:
            r += 1
        if r > max_row:
            break
        start_row = r
        while r <= max_row and row_non_empty_count(r) >= min_non_empty_per_row:
            r += 1
        end_row = r - 1

        min_c = None
        max_c = 0
        for rr in range(start_row, end_row+1):
            for cc in range(1, max_col+1):
                v = ws.cell(rr, cc).value
                if v is not None and str(v).strip() != "":
                    if min_c is None or cc < min_c:
                        min_c = cc
                    if cc > max_c:
                        max_c = cc
        if min_c is None:
            continue
        blocks.append((start_row, min_c, end_row, max_c))

    return blocks

def block_to_df(ws, block: Tuple[int,int,int,int]) -> pd.DataFrame:
    (min_row, min_col, max_row, max_col) = block
    data = []
    for r in ws.iter_rows(min_row=min_row, max_row=max_row, min_col=min_col, max_col=max_col, values_only=True):
        data.append(list(r))
    if not data:
        return pd.DataFrame()
    header = data[0]
    rows = data[1:]
    df = pd.DataFrame(rows, columns=header)
    df = df.dropna(how='all')
    return df

def get_images_with_positions(ws):
    """
    Return sorted list of (img, r1, c1, r2, c2) using 1-based rows/cols.
    Works for TwoCellAnchor and oneCell anchors.
    """
    out = []
    for img in getattr(ws, "_images", []):
        r1 = c1 = r2 = c2 = None
        a = getattr(img, "anchor", None)
        if a is not None:
            # Preferred: _from / _to
            if hasattr(a, "_from") and a._from:
                r1 = a._from.row + 1
                c1 = a._from.col + 1
            if hasattr(a, "_to") and a._to:
                r2 = a._to.row + 1
                c2 = a._to.col + 1
            # Fallback for attributes named 'to'
            if r2 is None and hasattr(a, "to") and a.to:
                r2 = a.to.row + 1
                c2 = a.to.col + 1
        # Fallbacks if nothing found:
        if r1 is None: r1 = 10**9
        if c1 is None: c1 = 10**9
        if r2 is None: r2 = r1
        if c2 is None: c2 = c1
        out.append((img, r1, c1, r2, c2))
    out.sort(key=lambda x: (x[1], x[2]))
    return out

def image_bytes(img) -> bytes:
    # openpyxl Image usually exposes a private _data() method that returns bytes
    for attr in ("_data",):
        if hasattr(img, attr):
            obj = getattr(img, attr)
            if callable(obj):
                try:
                    return obj()
                except Exception:
                    pass
    # As a last resort, see if there's a "path" (rare)
    if hasattr(img, "path") and img.path and os.path.exists(img.path):
        with open(img.path, "rb") as f:
            return f.read()
    raise RuntimeError("Unable to access image bytes from openpyxl image object.")

def scale_image_to_page(doc, image_stream: BytesIO, max_width_inches: float = 6.0):
    doc.add_picture(image_stream, width=Inches(max_width_inches))

def prepare_songs_and_images(excel_path: str, sheet_name: str = "Data"):
    wb = load_workbook(excel_path, data_only=True)
    ws = find_sheet_case_insensitive(wb, sheet_name)
    if ws is None:
        raise RuntimeError(f"Could not find a worksheet named '{sheet_name}'.")

    # Detect 4th block
    blocks = detect_blocks(ws, min_non_empty_per_row=1)
    if len(blocks) < 4:
        raise RuntimeError("Could not locate the fourth table/block on the 'Data' sheet.")
    target_block = blocks[3]
    (t_min_row, t_min_col, t_max_row, t_max_col) = target_block
    df = block_to_df(ws, target_block)

    # Find song title column
    song_col = None
    for col in df.columns:
        if isinstance(col, str) and any(k in col.lower() for k in ["titel", "song", "title", "lagu"]):
            song_col = col
            break
    if song_col is None:
        song_col = df.columns[0]

    # Build songs list
    songs = []
    for order_idx, val in enumerate(df[song_col].tolist()):
        if pd.isna(val):
            continue
        title = str(val).strip()
        if title:
            songs.append((order_idx, title))
    if not songs:
        raise RuntimeError("No song titles found in the fourth table/block.")

    # Song rows
    song_rows = [t_min_row + 1 + idx for idx, _ in songs]

    # Intervals
    intervals = [(song_rows[i], song_rows[i+1]) for i in range(len(song_rows)-1)]
    intervals.append((song_rows[-1], t_max_row + 1))

    # Get images
    all_imgs = get_images_with_positions(ws)
    imgs_in_block = [
        t for t in all_imgs
        if (t_min_row <= t[1] <= t_max_row) and (t_min_col <= t[2] <= t_max_col)
    ]

    from collections import defaultdict
    song_to_imgs = defaultdict(list)

    def relative_pos(r_abs, c_abs):
        return (r_abs - t_min_row, c_abs - t_min_col + 1)

    for (img, r1, c1, r2, c2) in imgs_in_block:
        idx = None
        for i, (start, end) in enumerate(intervals):
            if start <= r1 < end:
                idx = i
                break
        if idx is None:
            diffs = [abs(r1 - sr) for sr in song_rows]
            idx = diffs.index(min(diffs))
        title = songs[idx][1]
        r_rel, c_rel = relative_pos(r1, c1)
        song_to_imgs[title].append({
            "img": img,
            "abs_row": r1, "abs_col": c1,
            "rel_row": r_rel, "rel_col": c_rel,
        })

    return songs, song_rows, song_to_imgs

def _images_for_title(song_to_imgs: Dict[str, list], title: str):
    if title in song_to_imgs:
        return sorted(song_to_imgs[title], key=lambda d: (d["abs_col"], d["abs_row"]))
    # case-insensitive match as fallback
    tnorm = title.strip().casefold()
    for k in song_to_imgs.keys():
        if str(k).strip().casefold() == tnorm:
            return sorted(song_to_imgs[k], key=lambda d: (d["abs_col"], d["abs_row"]))
    return []

def add_song_sheet_entry(song_no: int, *, song_to_imgs: Dict[str, list], max_image_width_in: float = 6.0):
    if not (1 <= song_no <= len(songs)):
        raise ValueError(f"Song number {song_no} out of range (have {len(songs)}).")
    s = songs[song_no - 1]

    if s["Boek"] == "OPW":
        s["Boek"] = "Opwekking"

    # --- Heading (Roboto 11, bold, red) ---
    p = doc2.add_paragraph()
    pf = p.paragraph_format; pf.space_before = Pt(0); pf.space_after = Pt(0); pf.line_spacing = 1.0
    
    if is_empty(s["Nummer"]):
        r = p.add_run(f'{song_no}e lied: {s["Boek"]} "{s["Titel"]}"')
    else:
        if is_empty(s["Versen"]):
            r = p.add_run(f'{song_no}e lied: {s["Boek"]} {s["Nummer"]} "{s["Titel"]}"')
        else:
            r = p.add_run(f'{song_no}e lied: {s["Boek"]} {s["Nummer"]}: {s["Versen"]}  "{s["Titel"]}"')
        
    r.font.name = "Roboto"; r._element.rPr.rFonts.set(qn('w:eastAsia'), "Roboto"); r.font.size = Pt(11); r.bold = True
    from docx.shared import RGBColor; r.font.color.rgb = RGBColor(0xFF, 0x00, 0x00)

    # --- YouTube line (Roboto 11, bold, red) ---
    p = doc2.add_paragraph()
    pf = p.paragraph_format; pf.space_before = Pt(0); pf.space_after = Pt(0); pf.line_spacing = 1.0
    r = p.add_run(f'YouTube Reference: {s["youtube"]}')
    r.font.name = "Roboto"; r._element.rPr.rFonts.set(qn('w:eastAsia'), "Roboto"); r.font.size = Pt(11); r.bold = True
    r.font.color.rgb = RGBColor(0xFF, 0x00, 0x00)

    # --- one empty line (always) ---
    p = doc2.add_paragraph()
    pf = p.paragraph_format; pf.space_before = Pt(0); pf.space_after = Pt(0); pf.line_spacing = 1.0
    p.add_run("")

    # --- Images (if any) AFTER YouTube, then one empty line after the image(s) ---
    imgs_for_song = _images_for_title(song_to_imgs, str(s["Titel"]).strip())
    if imgs_for_song:
        for d in imgs_for_song:
            try:
                data = image_bytes(d["img"])
                bio = BytesIO(data)
                scale_image_to_page(doc2, bio, max_width_inches=max_image_width_in)
            except Exception as e:
                err_p = doc2.add_paragraph()
                err_pf = err_p.paragraph_format; err_pf.space_before = Pt(0); err_pf.space_after = Pt(0); err_pf.line_spacing = 1.0
                err_p.add_run(f"[Image could not be extracted: {e}]")
        # exactly one empty line after the (last) image
        p = doc2.add_paragraph()
        pf = p.paragraph_format; pf.space_before = Pt(0); pf.space_after = Pt(0); pf.line_spacing = 1.0
        p.add_run("")

    # --- Lyrics paragraphs (Roboto 11, black) ---
    for text in s["cells"]:
        p = doc2.add_paragraph()
        pf = p.paragraph_format; pf.space_before = Pt(0); pf.space_after = Pt(6); pf.line_spacing = 1.0
        r = p.add_run(text)
        r.font.name = "Roboto"; r._element.rPr.rFonts.set(qn('w:eastAsia'), "Roboto"); r.font.size = Pt(11); r.bold = False

    # --- page break at end of entry ---
    p = doc2.add_paragraph()
    p.add_run().add_break(WD_BREAK.PAGE)




# --- Set margins FIRST so header table can use full width ---
section = doc2.sections[0]
section.top_margin = Cm(0.5)
section.bottom_margin = Cm(1.2)
section.left_margin = Cm(1.5)
section.right_margin = Cm(1.5)

# Set a clean Normal style baseline
doc2.styles["Normal"].font.name = "Times New Roman"
doc2.styles["Normal"].font.size = Pt(11)

def add_para(text, bold=False, size=11, align=None, keep_next=False, space_before=0, space_after=3, line_spacing=1.0):
    p = doc2.add_paragraph()
    if align is not None:
        p.alignment = align
    run = p.add_run(text)
    run.bold = bool(bold)
    run.font.size = Pt(size)
    pf = p.paragraph_format
    pf.space_before = Pt(space_before)
    pf.space_after = Pt(space_after)
    pf.line_spacing = line_spacing
    if keep_next:
        pf.keep_with_next = True
    return p

# === HEADER TABLE (3 columns, 2 rows) ===
header = doc2.add_table(rows=2, cols=3)
header.autofit = False  # ensure widths stick

# Remove ALL borders first
for row in header.rows:
    for cell in row.cells:
        clear_cell_borders(cell)

# Stretch to full usable page width so long title fits on one line
usable_width = section.page_width - section.left_margin - section.right_margin
usable_w_in = usable_width / Inches(1)   # float inches

logo_w_in = 1.35
gap_w_in  = 0.08
text_w_in = max(4.5, usable_w_in - logo_w_in - gap_w_in)  # guard minimum

# Set column widths (both rows)
for r in (0, 1):
    set_cell_width(header.cell(r, 0), logo_w_in)
    set_cell_width(header.cell(r, 1), gap_w_in)
    set_cell_width(header.cell(r, 2), text_w_in)

# Merge first column rows and place logo
logo_cell = header.cell(0, 0).merge(header.cell(1, 0))
logo_par = logo_cell.paragraphs[0]
logo_par.alignment = WD_PARAGRAPH_ALIGNMENT.CENTER
logo_cell.vertical_alignment = WD_ALIGN_VERTICAL.CENTER
logo_run = logo_par.add_run()
logo_path = os.path.join(dir_path, LOGO_FILE)
if os.path.exists(logo_path):
    logo_run.add_picture(logo_path, width=Inches(min(1.2, logo_w_in - 0.1)))
else:
    logo_par.add_run("[LOGO ontbreekt: logo.png]").bold = True

# Column 2 (gap) stays empty and borderless
header.cell(0, 1).text = ""
header.cell(1, 1).text = ""

# Column 3, Row 1: titles, tight spacing; ONLY bottom double border
title_cell_top = header.cell(0, 2); title_cell_top.text = ""

p1 = title_cell_top.paragraphs[0]
p1.alignment = WD_PARAGRAPH_ALIGNMENT.LEFT
p1.paragraph_format.space_after = Pt(0)
p1.paragraph_format.line_spacing = 1.0
r1 = p1.add_run("GEREJA KRISTEN INDONESIA NEDERLAND")
r1.bold = True; r1.font.size = Pt(16); r1.font.name = "Bookman Old Style"

p2 = title_cell_top.add_paragraph()
p2.alignment = WD_PARAGRAPH_ALIGNMENT.LEFT
p2.paragraph_format.space_after = Pt(2)
p2.paragraph_format.line_spacing = 1.0
r2 = p2.add_run("Indonesisch Nederlands Christelijke Kerk")
r2.bold = True; r2.italic = True; r2.font.size = Pt(11); r2.font.name = "Times New Roman"

# ONLY border on this cell
set_cell_borders(title_cell_top, bottom=("double",12))

# Column 3, Row 2: remove default blank paragraph, then nested 3-col table (borderless)
title_cell_bot = header.cell(1, 2)

# Remove auto empty paragraph so there's NO blank line
while title_cell_bot.paragraphs:
    p = title_cell_bot.paragraphs[0]._element
    p.getparent().remove(p)

nested = title_cell_bot.add_table(rows=3, cols=3)
nested.autofit = False
for r in nested.rows:
    for c in r.cells:
        clear_cell_borders(c)

# widths inside nested: label 2.25", ":" 0.25", value = remainder
label_w, colon_w = 2.25, 0.25
value_w = max(1.0, text_w_in - label_w - colon_w)
for i in range(3):
    set_cell_width(nested.cell(i,0), label_w)
    set_cell_width(nested.cell(i,1), colon_w)
    set_cell_width(nested.cell(i,2), value_w)

labels = ["Eredienst van zondag", "Voorganger", "Ouderling van Dienst"]
values = [long_date, voorganger, ouderling]

for i in range(3):
    lp = nested.cell(i,0).paragraphs[0]
    lp.paragraph_format.space_after = Pt(0)
    lp.paragraph_format.line_spacing = 1.0
    lrun = lp.add_run(labels[i])
    lrun.font.size = Pt(12); lrun.font.name = "Times New Roman"; lrun.bold = True

    cp = nested.cell(i,1).paragraphs[0]
    cp.alignment = WD_PARAGRAPH_ALIGNMENT.CENTER
    cp.paragraph_format.space_after = Pt(0)
    cp.paragraph_format.line_spacing = 1.0
    crun = cp.add_run(":")
    crun.font.size = Pt(12); crun.font.name = "Times New Roman"; crun.bold = True

    vp = nested.cell(i,2).paragraphs[0]
    vp.paragraph_format.space_after = Pt(0)
    vp.paragraph_format.line_spacing = 1.0
    vrun = vp.add_run(values[i])
    vrun.font.size = Pt(12); vrun.font.name = "Times New Roman"; vrun.bold = True

p = doc2.add_paragraph()
pf = p.paragraph_format
pf.space_before = Pt(0); pf.space_after = Pt(6); pf.line_spacing = 1.0; pf.space_before = Pt(6)
r = p.add_run("")
r.font.name = "Roboto"; r.font.size = Pt(11)

# --- List all songs, one per line ---
counter = 0
for s in songs:
    if s["Boek"] == "Kosong":
        buku = ""
    else:
        buku = s["Boek"]
    
    counter = counter + 1
    p = doc2.add_paragraph()
    pf = p.paragraph_format
    pf.space_before = Pt(0)
    pf.space_after  = Pt(0)
    pf.line_spacing = 1.0

    if is_empty(s["Nummer"]):
        run = p.add_run(f'{counter}e lied: {buku} "{s["Titel"]}"')
    elif is_empty(s["Versen"]):
        run = p.add_run(f'{counter}e lied: {buku} {s["Nummer"]} "{s["Titel"]}"')
    else:
        run = p.add_run(f'{counter}e lied: {buku} {s["Nummer"]}: {s["Versen"]}  "{s["Titel"]}"')
    run.font.name = "Roboto"
    run.font.size = Pt(12)
    run.bold = True

p = doc2.add_paragraph()
pf = p.paragraph_format
pf.space_before = Pt(0); pf.space_after = Pt(6); pf.line_spacing = 1.0; pf.space_before = Pt(6)
r = p.add_run("")
r.font.name = "Roboto"; r.font.size = Pt(11)
_, _, song_to_imgs = prepare_songs_and_images(os.path.join(dir_path + '/file mingguan/', EXCEL_NAME))
                                        
add_song_sheet_entry(1, song_to_imgs=song_to_imgs, max_image_width_in=6.0)
add_song_sheet_entry(2, song_to_imgs=song_to_imgs, max_image_width_in=6.0)
add_song_sheet_entry(3, song_to_imgs=song_to_imgs, max_image_width_in=6.0)
add_song_sheet_entry(4, song_to_imgs=song_to_imgs, max_image_width_in=6.0)
add_song_sheet_entry(5, song_to_imgs=song_to_imgs, max_image_width_in=6.0)
add_song_sheet_entry(6, song_to_imgs=song_to_imgs, max_image_width_in=6.0)
add_song_sheet_entry(7, song_to_imgs=song_to_imgs, max_image_width_in=6.0)





# -------------------------------
# Save
# -------------------------------
out_name = f"LiturgieB {short_date}_{REGIO}.docx"
out_path = os.path.join(dir_path + '/file mingguan/', out_name)
doc2.save(out_path)
print(f"Word document generated: {out_path}")




# -------------------------------
# PPT
# -------------------------------

from pptx import Presentation
from pptx.util import Emu, Inches, Pt, Cm
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE, MSO_ANCHOR
from pptx.oxml.xmlchemy import OxmlElement
from pptx.oxml.ns import qn
import textwrap
from lxml import etree as ET

# === Slide geometry (from your template) ===
SLIDE_WIDTH = Emu(12192000)   # 13.333 in (16:9)
SLIDE_HEIGHT = Emu(6858000)   # 7.5 in
BOX_TOP = Cm(0.0)
BOX_LEFT = Cm(8.5)
BOX_WIDTH = Cm(17)
yellow = RGBColor(255,173,3)
white = RGBColor(255, 255, 255)


# --- Config you can tweak ---
BG_HEX = "1E1947"                 # slide background
BOX_W = Cm(17)                    # 17 cm wide
BOX_H = Cm(17)                    # 17 cm high
FONT_NAME = "Calibri"
FONT_SIZE_PT = 32
VERSE_NUM_COLOR = RGBColor(255, 173, 3)  # (255,173,3)
TEXT_COLOR = RGBColor(255, 255, 255)

# Heuristic for pagination (no font metrics in python-pptx).
# With Calibri 32pt in a 17×17 cm box, ~13 lines fits comfortably at line_spacing≈1.0.
MAX_LINES_PER_SLIDE = 10
# Rough characters per wrapped line (tune if needed; affects how many verses fit per slide):
CHARS_PER_LINE = 30

def set_background(slide, hex_rgb="1E1947"):
    """Set slide background to solid color."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor.from_string(hex_rgb)

def add_textbox(slide, left, top, width, height, text, font_name, font_size=36, italics=True, bold=True, alignment="CENTER", color = RGBColor(255,255,255)):
    """Add a text box with specific formatting."""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True                  # wrap long lines
    tf.auto_size = MSO_AUTO_SIZE.NONE    # keep the box size fixed (don’t auto-resize)
    tf.vertical_anchor = MSO_ANCHOR.TOP  # start from top of the box
    p = tf.paragraphs[0]
    p.alignment = getattr(PP_ALIGN, alignment)
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italics
    run.font.color.rgb = color
    return box

def set_tabs(paragraph, stops_cm):
    pPr = paragraph._p.get_or_add_pPr()
    tabLst = OxmlElement('a:tabLst')
    for stop in stops_cm:
        t = OxmlElement('a:tab')
        t.set('pos', str(Cm(stop).emu))  # tab position
        t.set('algn', 'l')               # left-aligned tab
        tabLst.append(t)
    pPr.append(tabLst)


def _to_int_or_raise(x, field):
    try:
        return int(str(x).strip())
    except Exception:
        raise ValueError(f"'{field}' moet een geheel getal zijn, gekregen: {x!r}")

def _flatten_text(node):
    """Flatten the mixed JSON 'content' into plain text."""
    if node is None:
        return ""
    if isinstance(node, str):
        return node
    if isinstance(node, list):
        return "".join(_flatten_text(x) for x in node)
    if isinstance(node, dict):
        if "text" in node and isinstance(node["text"], str):
            return node["text"]
        t = node.get("type")
        if t in ("line-break", "softbreak", "hardbreak"):
            return " "  # or "\n" if you prefer hard line breaks
        if "content" in node and isinstance(node["content"], list):
            return "".join(_flatten_text(x) for x in node["content"])
    return ""

def _collect_verses_from_chapter(chapter_content, v_from, v_to):
    """
    chapter_content: list of paragraphs, each with 'content' items.
    Collect items where type=='verse-text' and parse verseId like 'GEN.1.18'.
    Returns {verse_num: text} for the inclusive range.
    """
    verses = {v: "" for v in range(v_from, v_to + 1)}
    for paragraph in chapter_content:
        for item in paragraph.get("content", []):
            if item.get("type") == "verse-text":
                verse_id = item.get("verseId")  # e.g. 'GEN.1.18'
                try:
                    verse_num = int(str(verse_id).split(".")[-1])
                except Exception:
                    continue
                if v_from <= verse_num <= v_to:
                    verses[verse_num] += _flatten_text(item.get("content"))

    # normalize whitespace
    for v in list(verses.keys()):
        verses[v] = " ".join(verses[v].split())
    return verses

def _estimate_wrapped_lines(text: str, chars_per_line: int = CHARS_PER_LINE) -> int:
    """Heuristic: wrap by character count to estimate how many lines the verse will occupy."""
    # Replace multiple spaces/newlines with single spaces for a stable estimate
    norm = re.sub(r"\s+", " ", text.strip())
    wrapped = textwrap.wrap(norm, width=chars_per_line, break_long_words=False, break_on_hyphens=True)
    return max(1, len(wrapped))

def _new_slide_with_box(prs):
    """Create a blank slide with background and a centered 17×17 cm text box (vertical middle, left-aligned)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    # background
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor.from_string(BG_HEX)
    # center the 17×17 box on the slide
    left = (prs.slide_width - BOX_W) / 2
    top = (prs.slide_height - BOX_H) / 2
    box = slide.shapes.add_textbox(left, top, BOX_W, BOX_H)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE   # vertically centered
    # margins small but breathable
    tf.margin_left = Cm(0.3)
    tf.margin_right = Cm(0.3)
    tf.margin_top = Cm(0.2)
    tf.margin_bottom = Cm(0.2)
    return slide, tf

def add_verses_to_ppt(prs: Presentation, dutch_book_name: str, chapter_number, verse_from, verse_to=None,
                      font_name: str = FONT_NAME, font_size_pt: int = FONT_SIZE_PT):
    """
    - Maps Dutch book name -> abbreviation
    - Opens JSON file named '{ABBR}.{chapter}'
    - Gets verses in inclusive range [verse_from, verse_to]
      * If verse_to is None/""/0 -> uses single verse (verse_to = verse_from)
    - Adds one or more slides to 'prs'.
      * Each slide has a 17×17 cm box, Calibri 32, white text, left-aligned, vertically centered.
      * Pagination: we never split a verse across slides. If the next verse doesn't fit heuristically,
        we start a new slide.
      * Each verse begins with its number as superscript, size 32, bold, color (255,173,3).
      * ▶ NEW: Adds a top-right label "<Book> <Chapter> (NBV)" in yellow on every slide.
    """
    # Normalize inputs
    chapter_number = _to_int_or_raise(chapter_number, "hoofdstuk")
    verse_from = _to_int_or_raise(verse_from, "vers_van")
    if verse_to in (None, "", 0, "0"):
        verse_to = verse_from
    else:
        verse_to = _to_int_or_raise(verse_to, "vers_tot")
    if verse_to < verse_from:
        verse_from, verse_to = verse_to, verse_from

    # Map book to abbreviation
    try:
        abbr = bible_books_nl[dutch_book_name]
    except KeyError:
        raise ValueError(f"Onbekende boeknaam: {dutch_book_name!r}")

    # Load JSON chapter
    json_path = os.path.join(dir_path, "bible", f"{abbr}.{chapter_number}")
    if not os.path.exists(json_path):
        # Try downloading from Dropbox on demand (web deployment)
        _fetch = globals().get('_ensure_bible_file')
        if _fetch:
            try:
                _fetch(f"{abbr}.{chapter_number}")
            except Exception:
                pass
    if not os.path.exists(json_path):
        # fall back to .json if your files use that extension
        if os.path.exists(json_path + ".json"):
            json_path = json_path + ".json"
        else:
            raise FileNotFoundError(f"Bestand niet gevonden: {json_path}")

    with open(json_path, "r", encoding="utf-8") as f:
        payload = json.load(f)

    try:
        chapter = payload["data"]["chapter"]
        chapter_content = chapter["content"]
    except Exception as e:
        raise ValueError(f"Onverwachte JSON-structuur in {json_path}: {e}")

    verses = _collect_verses_from_chapter(chapter_content, verse_from, verse_to)
    if not any(verses.values()):
        return prs  # nothing to add

    # Helper: add the yellow top-right label on a given slide
    def _add_top_right_label(slide):
        label_text = f"{dutch_book_name} {chapter_number} (NBV 21)"
        label_w = BOX_WIDTH
        label_h = Cm(2.0)
        label_left = BOX_LEFT   # 1 cm right margin
        label_top = Cm(0.5)                                # 0.5 cm from top
        box = slide.shapes.add_textbox(label_left, label_top, label_w, label_h)
        tf = box.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        r = p.add_run()
        r.text = label_text
        r.font.name = font_name
        r.font.size = Pt(24)                # small header size
        r.font.bold = True                  # matches your style
        r.font.color.rgb = VERSE_NUM_COLOR  # yellow (255,173,3)

    # Build slides with pagination by verse
    lines_used = 0
    slide, tf = _new_slide_with_box(prs)
    _add_top_right_label(slide)

    for v in range(verse_from, verse_to + 1):
        txt = verses.get(v, "")
        txt = normalize_spaces(txt)
        
        if not txt:
            continue

        # Estimate how many lines this verse will take (roughly)
        est_lines = _estimate_wrapped_lines(txt, CHARS_PER_LINE)

        # If it won't fit on current slide, start a new one
        if lines_used > 0 and (lines_used + est_lines) > MAX_LINES_PER_SLIDE:
            lines_used = 0
            slide, tf = _new_slide_with_box(prs)
            _add_top_right_label(slide)

        # Add the verse as a paragraph (number as superscript, then text)
        p = tf.paragraphs[0] if lines_used == 0 and len(tf.paragraphs) == 1 and tf.paragraphs[0].text == "" else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT

        # verse number (superscript, bold, yellow)
        r_num = p.add_run()
        r_num.text = str(v)
        r_num.font.name = font_name
        r_num.font.size = Pt(font_size_pt)
        r_num.font.bold = True
        r_num.font.color.rgb = VERSE_NUM_COLOR
        r_num.font.superscript = True
        r_num._r.get_or_add_rPr().set("baseline", "30000") # enforce superscript

        # space and verse text (white)
        r_txt = p.add_run()
        r_txt.text = " " + txt
        r_txt.font.name = font_name
        r_txt.font.size = Pt(font_size_pt)
        r_txt.font.bold = False
        r_txt.font.color.rgb = TEXT_COLOR
        r_txt.font.superscript = False

        # Update our estimate of used lines
        lines_used += est_lines

    p = tf.add_paragraph()
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "\n-"
    r.font.name = font_name
    r.font.size = Pt(font_size_pt)
    r.font.bold = False
    r.font.color.rgb = TEXT_COLOR


    return prs

# --- Indonesian Bible to PPT (TB / TB2, Zefania XML, tolerant parser) ---
# Re-use your existing constants/helpers already in your file:
# BG_HEX, BOX_W, BOX_H, FONT_NAME, FONT_SIZE_PT, VERSE_NUM_COLOR, TEXT_COLOR,
# MAX_LINES_PER_SLIDE, CHARS_PER_LINE, _to_int_or_raise, _estimate_wrapped_lines, _new_slide_with_box

# 1) Book sets to choose TB (OT) vs TB2 (NT)
OT_ID_BOOKS = {
    "Kejadian","Keluaran","Imamat","Bilangan","Ulangan","Yosua","Hakim-hakim","Rut",
    "1 Samuel","2 Samuel","1 Raja-raja","2 Raja-raja","1 Tawarikh","2 Tawarikh","Ezra",
    "Nehemia","Ester","Ayub","Mazmur","Amsal","Pengkhotbah","Kidung Agung","Yesaya",
    "Yeremia","Ratapan","Yehezkiel","Daniel","Hosea","Yoel","Amos","Obaja","Yunus",
    "Mikha","Nahum","Habakuk","Zefanya","Hagai","Zakharia","Maleakhi"
}
NT_ID_BOOKS = {
    "Matius","Markus","Lukas","Yohanes","Kisah Para Rasul","Roma","1 Korintus","2 Korintus",
    "Galatia","Efesus","Filipi","Kolose","1 Tesalonika","2 Tesalonika","1 Timotius","2 Timotius",
    "Titus","Filemon","Ibrani","Yakobus","1 Petrus","2 Petrus","1 Yohanes","2 Yohanes",
    "3 Yohanes","Yudas","Wahyu"
}



# Aliases so inputs like "Act", "KPR", "Kisah" map to the exact bname used in the XML
BOOK_ALIASES_ID = {
    "kisah para rasul": "Kisah Para Rasul",
    "act": "Kisah Para Rasul",
    "acts": "Kisah Para Rasul",
    "kpr": "Kisah Para Rasul",
    # add other handy aliases if you need
}

def _norm(s: str) -> str:
    return " ".join(
        (s or "")
        .strip()
        .lower()
        .replace("’", "'").replace("‐", "-").replace("–", "-")
        .split()
    )

def _resolve_indonesian_book_name(raw_name: str, root: ET._Element) -> str:
    """
    Resolve user input to an exact @bname present in the XML.
    - Try alias map
    - Try exact (normalized) equality over all BIBLEBOOK/@bname
    Raise if not found.
    """
    want_norm = _norm(raw_name)
    # alias first
    if want_norm in BOOK_ALIASES_ID:
        target = BOOK_ALIASES_ID[want_norm]
        return target

    # exact normalized match across available bnames
    candidates = {}
    for bb in root.findall(".//BIBLEBOOK"):
        bname = bb.get("bname") or ""
        candidates[_norm(bname)] = bname

    if want_norm in candidates:
        return candidates[want_norm]

    # last resort: fail clearly (avoid substring matches that can mislead)
    raise ValueError(f"Buku tidak ditemukan di XML: {raw_name!r}. Tersedia: {sorted(set(candidates.values()))[:10]} ...")

def _get_chapter_verses_from_xml(xml_path: str, indo_book_name: str, chapter_number: int,
                                 v_from: int, v_to: int) -> dict[int, str]:
    """
    Parse Zefania XML and return {verse_num -> text} for the requested range.
    Structure:
      <XMLBIBLE>
        <BIBLEBOOK bname="Kisah Para Rasul">
          <CHAPTER cnumber="2">
            <VERS vnumber="17"> ... </VERS>
          </CHAPTER>
        </BIBLEBOOK>
      </XMLBIBLE>
    """
    if not os.path.exists(xml_path):
        raise FileNotFoundError(f"Tidak ditemukan: {xml_path}")

    parser = ET.XMLParser(recover=True)
    tree = ET.parse(xml_path, parser=parser)
    root = tree.getroot()

    # normalize
    want_book = _norm(indo_book_name)

    # Find exact book element
    book_el = None
    for bb in root.findall(".//BIBLEBOOK"):
        bname = _norm(bb.get("bname"))
        if bname == want_book:
            book_el = bb
            break
    if book_el is None:
        raise ValueError(f"Buku {indo_book_name!r} tidak ditemukan di {xml_path}")

    # Find exact chapter element (integer equality, no substrings!)
    chap_el = None
    for ch in book_el.findall("./CHAPTER"):
        try:
            if int(ch.get("cnumber")) == int(chapter_number):
                chap_el = ch
                break
        except Exception:
            continue
    if chap_el is None: 
        raise ValueError(f"Pasal {chapter_number} tidak ditemukan untuk {indo_book_name!r}")

    # Collect verses
    out = {}
    for vs in chap_el.findall("./VERS"):
        try:
            vn = int(vs.get("vnumber"))
        except Exception:
            continue
        if v_from <= vn <= v_to:
            txt = "".join(vs.itertext())
            out[vn] = " ".join((txt or "").split())

    return out

def add_verses_to_ppt_indo(
    prs,
    indo_book_name: str,
    chapter_number,
    verse_from,
    verse_to=None,
    dir_path=dir_path,
    font_name="Calibri",
    font_size_pt: int = 32,
):
    """
    Indonesian TB/TB2 (Zefania XML):
    - Picks TB for OT, TB2 for NT, from <dir_path>/bible/.
    - Adds one or more slides with centered 17×17 cm text box (LEFT align, vertical MIDDLE).
    - Calibri 32, white; verse numbers are superscript, bold, yellow (255,173,3).
    - Top-right yellow header: "<Book> <Chapter:From–To> (TB|TB2)" on every slide.
    Returns number of slides added.
    """
    # normalize range
    chapter_number = int(chapter_number)
    verse_from = int(verse_from)
    verse_to = verse_from if verse_to in (None, "", 0, "0") else int(verse_to)
    if verse_to < verse_from:
        verse_from, verse_to = verse_to, verse_from

    # pick file & label
    bible_dir = os.path.join(dir_path, "bible")
    if not os.path.isdir(bible_dir):
        raise FileNotFoundError(f"Folder Alkitab tidak ditemukan: {bible_dir}")
    # decide OT vs NT by resolved exact bname
    # (we need the XML open to resolve; open TB first just to get bnames)
    # Simpler: use your existing sets
    label = "TB2" if _norm(indo_book_name) in {_norm(x) for x in NT_ID_BOOKS} else "TB"
    xml_file = "Alkitab_PB_TB2.xml" if label == "TB2" else "Alkitab_TB.xml"
    xml_path = os.path.join(bible_dir, xml_file)
    if not os.path.exists(xml_path):
        # Try downloading from Dropbox on demand (web deployment)
        _fetch = globals().get('_ensure_bible_file')
        if _fetch:
            try:
                _fetch(xml_file)
            except Exception:
                pass

    # read verses strictly from the requested chapter
    verses = _get_chapter_verses_from_xml(xml_path, indo_book_name, chapter_number, verse_from, verse_to)
    if not verses:
        raise ValueError(f"Tidak ada ayat ditemukan untuk {indo_book_name} {chapter_number}:{verse_from}-{verse_to} di {label}")

    # Build the header "<Book> <Chapter:From–To> (TB|TB2)"
    header_text = f"{indo_book_name} {chapter_number} ({label})"

    # helper to add header on a slide (top-right)
    def _add_top_right_label(slide):
        w, h = BOX_WIDTH, Cm(2.0)
        left = BOX_LEFT
        top = Cm(0.5)
        box = slide.shapes.add_textbox(left, top, w, h)
        tf = box.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        r = p.add_run()
        r.text = header_text
        r.font.name = font_name
        r.font.size = Pt(24)
        r.font.bold = True
        r.font.color.rgb = RGBColor(255, 173, 3)

    # build slides (reuse your _new_slide_with_box / _estimate_wrapped_lines constants)
    slides_added = 0
    lines_used = 0
    slide, tf = _new_slide_with_box(prs)
    _add_top_right_label(slide)
    slides_added += 1

    for vn in range(verse_from, verse_to + 1):
        txt = verses.get(vn, "")
        txt = normalize_spaces(txt)
        if not txt:
            continue

        est_lines = _estimate_wrapped_lines(txt, CHARS_PER_LINE)
        if lines_used > 0 and (lines_used + est_lines) > MAX_LINES_PER_SLIDE:
            lines_used = 0
            slide, tf = _new_slide_with_box(prs)
            _add_top_right_label(slide)
            slides_added += 1

        # paragraph
        p = tf.paragraphs[0] if (lines_used == 0 and len(tf.paragraphs) == 1 and tf.paragraphs[0].text == "") else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT

        # superscript "chapter:verse" to remove any ambiguity like 17:18
        r_num = p.add_run()
        r_num.text = f"{vn}"
        r_num.font.name = font_name
        r_num.font.size = Pt(font_size_pt)
        r_num.font.bold = True
        r_num.font.color.rgb = RGBColor(255, 173, 3)
        r_num.font.superscript = True
        r_num._r.get_or_add_rPr().set("baseline", "30000")

        r_txt = p.add_run()
        r_txt.text = " " + txt
        r_txt.font.name = font_name
        r_txt.font.size = Pt(font_size_pt)
        r_txt.font.bold = False
        r_txt.font.color.rgb = RGBColor(255, 255, 255)
        lines_used += est_lines

    
    p = tf.add_paragraph()
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "\n-"
    r.font.name = font_name
    r.font.size = Pt(font_size_pt)
    r.font.bold = False
    r.font.color.rgb = RGBColor(255, 255, 255)
    
    return slides_added

def add_sermon_doc_to_ppt(
    prs,
    docx_path,
    font_name="Calibri",
    font_size_pt=32,
    box_width=Cm(17),
    box_height=Cm(17),
    box_left=Cm(2),
    box_top=Cm(0),
    bg_hex="1E1947",
    max_lines_per_slide=None,
    chars_per_line=None,
):
    """
    - Preserve Word paragraphing (no auto-merging).
    - Add ONE blank line between paragraphs when multiple fit on the same slide.
    - Prefer one paragraph per slide; if it doesn't fit with current content but fits alone,
      start a NEW slide and place it intact (no splitting).
    - ONLY when a paragraph itself is too long for a slide, split it into CHUNKS of sentences
      that fit. Each chunk is rendered as ONE PPT paragraph (wrapping handled by PowerPoint).
    - 17×17 cm textbox (defaults), left-aligned, vertical middle, Calibri 32, white.
    """

    # Heuristics (fallbacks if not provided elsewhere)
    try:
        ML = max_lines_per_slide or MAX_LINES_PER_SLIDE
    except NameError:
        ML = max_lines_per_slide or 13

    try:
        CPL = chars_per_line or CHARS_PER_LINE
    except NameError:
        CPL = chars_per_line or 30

    def _new_slide_and_textframe():
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor.from_string(bg_hex)

        box = slide.shapes.add_textbox(box_left, box_top, box_width, box_height)
        tf = box.text_frame
        tf.clear()
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.NONE
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = tf.margin_right = Cm(0.3)
        tf.margin_top = tf.margin_bottom = Cm(0.2)
        return slide, tf

    def _wrap_lines(s: str) -> list[str]:
        norm = re.sub(r"\s+", " ", (s or "").strip())
        lines = textwrap.wrap(norm, width=CPL, break_long_words=False, break_on_hyphens=True)
        return lines or [""]

    def _estimate_lines(s: str) -> int:
        return max(1, len(_wrap_lines(s)))

    def _sentences(s: str) -> list[str]:
        """
        Split into sentences; keep punctuation. Handles ., ?, ! (optionally followed by closing quote/bracket).
        """
        text = " ".join((s or "").split())
        if not text:
            return []
        parts = re.findall(r'.+?(?:[.?!][\'")\]]?(?=\s|$)|$)', text)
        return [p.strip() for p in parts if p.strip()]

    def _add_blank_line(tf):
        p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run(); r.text = ""
        r.font.name = font_name
        r.font.size = Pt(font_size_pt)
        r.font.color.rgb = RGBColor(255, 255, 255)

    def _add_one_paragraph(tf, text):
        """Render a paragraph/chunk as ONE PPT paragraph; let word_wrap handle visual lines."""
        if len(tf.paragraphs) == 1 and tf.paragraphs[0].text == "":
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = " ".join((text or "").split())
        r.font.name = font_name
        r.font.size = Pt(font_size_pt)
        r.font.color.rgb = RGBColor(255, 255, 255)

    # --- Load Word document; preserve paragraphs as-is ---
    doc = Document(docx_path)
    paragraphs = [p.text.strip() for p in doc.paragraphs if (p.text or "").strip() != ""]
    if not paragraphs:
        return 0

    slides_added = 0
    slide, tf = _new_slide_and_textframe()
    slides_added += 1
    used_lines = 0

    for para in paragraphs:
        if not para:
            continue

        para_lines = _estimate_lines(para)
        spacer_cost = 1 if used_lines > 0 else 0  # blank line between paragraphs on same slide

        if used_lines + spacer_cost + para_lines <= ML:
            # 1) Fits on current slide → pack it
            if spacer_cost:
                _add_blank_line(tf)
                used_lines += 1
            _add_one_paragraph(tf, para)
            used_lines += para_lines
            continue

        if para_lines <= ML:
            # 2) Doesn't fit with current content, but fits alone → new slide, no splitting
            slide, tf = _new_slide_and_textframe()
            slides_added += 1
            used_lines = 0
            _add_one_paragraph(tf, para)
            used_lines = para_lines
            continue

        # 3) Paragraph itself is longer than a slide → split into sentence CHUNKS that fit
        sents = _sentences(para)
        idx = 0
        while idx < len(sents):
            # if there is already content, move to a fresh slide
            if used_lines > 0:
                slide, tf = _new_slide_and_textframe()
                slides_added += 1
                used_lines = 0

            # Build a chunk by adding sentences until adding the next would exceed ML
            chunk_text = ""
            chunk_lines = 0
            made_progress = False

            while idx < len(sents):
                candidate = (chunk_text + " " + sents[idx]).strip() if chunk_text else sents[idx]
                cand_lines = _estimate_lines(candidate)

                # If chunk already has content and adding this sentence would overflow, stop chunk here
                if chunk_text and cand_lines > ML:
                    break

                # If one single sentence itself exceeds ML: place it alone on this slide (no splitting)
                if not chunk_text and _estimate_lines(sents[idx]) > ML:
                    chunk_text = sents[idx]
                    chunk_lines = _estimate_lines(chunk_text)
                    idx += 1
                    made_progress = True
                    break

                # Otherwise accept this sentence into the chunk
                chunk_text = candidate
                chunk_lines = cand_lines
                idx += 1
                made_progress = True

            # Render this chunk as ONE PPT paragraph
            _add_one_paragraph(tf, chunk_text)
            used_lines = chunk_lines

            # If more sentences remain, go to a fresh slide
            if idx < len(sents):
                slide, tf = _new_slide_and_textframe()
                slides_added += 1
                used_lines = 0

            # Safety: if we didn't consume anything (shouldn't happen), break to avoid loops
            if not made_progress:
                break

    return slides_added


# === MAIN ===
prs = Presentation()
prs.slide_width = SLIDE_WIDTH
prs.slide_height = SLIDE_HEIGHT


# --- INTRO ---


# --- Slide 1 ---
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_background(slide1)

# Dutch notice (top)
add_textbox(
    slide1,
    BOX_LEFT, BOX_TOP, BOX_WIDTH, Cm(8),
    "Gaarne uw mobiele telefoon gedurende de eredienst uitzetten of op “stil” zetten. Dank u.",
    "Calibri", 36, italics=False
)


# Center logo (optional placeholder, if you have a logo.png)
phone_path = os.path.join(dir_path,"telephone.gif")
if os.path.exists(logo_path):
    slide1.shapes.add_picture(phone_path, Cm(15), Cm(7.5), width=Cm(4.5), height=Cm(4.5))

# Indonesian notice (bottom)
add_textbox(
    slide1,
    BOX_LEFT, Cm(12), BOX_WIDTH, Cm(4.5),
    "Mohon hand-phone anda dimatikan selama kebaktian berlangsung. Terima kasih.",
    "Garamond", italics=True
)

# --- Slide 2 ---
slide2 = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_background(slide2)
add_textbox(
    slide2,
    BOX_LEFT, BOX_TOP, BOX_WIDTH, Cm(2), "GKIN Eredienst", "Calibri", 54, False
)

add_textbox(
    slide2,
    BOX_LEFT, Cm(1.8), BOX_WIDTH, Cm(2), "Online Landelijke Eredienst", "Calibri", 32, False
)

add_textbox(
    slide2,
    BOX_LEFT, Cm(3.0), BOX_WIDTH, Cm(2), long_date, "Calibri", 32, False
)

# add pictures
line_path = os.path.join(dir_path,"garis.png")
if os.path.exists(line_path):
    slide2.shapes.add_picture(line_path, BOX_LEFT, Cm(4.5), BOX_WIDTH, Cm(0.5))

logo2_path = os.path.join(dir_path,"logo2.png")
if os.path.exists(logo2_path):
    slide2.shapes.add_picture(logo2_path, Cm(23), Cm(16), Cm(2.47), Cm(2.47))

# add table
box = slide2.shapes.add_textbox(BOX_LEFT, Cm(5.1), BOX_WIDTH, Cm(16))
tf = box.text_frame
tf.word_wrap = True
tf.clear()

lines = [
    f"Voorganger\t:\t{voorganger}",
    f"Ouderling v. dienst\t:\t{ouderling}",
    f"Eerste Ontvangst\t:\t{eersteO} ",
]

leden = re.split(r'[;,]\s*', muzikanten)
for i, lid in enumerate(leden):
    if i == 0:
        lines.append(f"Muzikanten\t:\t{lid}")
    else:
        lines.append(f"\t\t{lid.strip()}")    

leden = re.split(r'[;,]\s*', voorzangers)
for i, lid in enumerate(leden):
    if i == 0:
        lines.append(f"Voorzangers\t:\t{lid}")
    else:
        lines.append(f"\t\t{lid.strip()}")

leden = re.split(r'[;,]\s*', beamer)
for i, lid in enumerate(leden):
    if i == 0:
        lines.append(f"Beamer\t:\t{lid}")
    else:
        lines.append(f"\t\t{lid.strip()}")

leden = re.split(r'[;,]\s*', geluid)
for i, lid in enumerate(leden):
    if i == 0:
        lines.append(f"Geluid\t:\t{lid}")
    else:
        lines.append(f"\t\t{lid.strip()}")

leden = re.split(r'[;,]\s*', knd)
for i, lid in enumerate(leden):
    if i == 0:
        lines.append(f"KND\t:\t{lid}")
    else:
        lines.append(f"\t\t{lid.strip()}")

leden = re.split(r'[;,]\s*', tieners)
if is_empty(leden[0]) == False:
    for i, lid in enumerate(leden):
        if i == 0:
            lines.append(f"Tieners\t:\t{lid}")
        else:
            lines.append(f"\t\t{lid.strip()}")

counter = 0

for line in lines:
    if counter == 0:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    counter = counter + 1
    set_tabs(p, [6.5, 7.0])
    p.text = line
    p.font.size = Pt(18)
    p.font.name = "Calibri"
    p.font.color.rgb = white


# --- MEDEDELINGEN ---
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_background(slide)
box = slide.shapes.add_textbox(BOX_LEFT, Cm(1), BOX_WIDTH, Cm(16))
tf = box.text_frame
tf.clear()
tf.word_wrap = True
tf.auto_size = MSO_AUTO_SIZE.NONE
tf.vertical_anchor = MSO_ANCHOR.MIDDLE

p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = "MEDEDELINGEN\n"
r.font.name = "Calibri"
r.font.size = Pt(48)
r.font.color.rgb = white
r.font.bold = True

p = tf.add_paragraph()
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = "PEMBERITAHUAN" 
r.font.name = "Calibri"
r.font.size = Pt(48)
r.font.color.rgb = white
r.font.italic = True
r.font.bold = True

# --- woord van welkom ---

slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_background(slide)
box = slide.shapes.add_textbox(BOX_LEFT, Cm(1), BOX_WIDTH, Cm(16))
tf = box.text_frame
tf.clear()
tf.word_wrap = True
tf.auto_size = MSO_AUTO_SIZE.NONE
tf.vertical_anchor = MSO_ANCHOR.MIDDLE

p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = "Selamat pagi saudara-saudari, Atas nama dewan Majelis regio Amstelveen, saya mengucapkan selamat datang kepada saudara-saudari sekalian dalam kebaktian ini, khususnya bagi mereka yang baru pertama kali hadir."
r.font.name = "Calibri"
r.font.size = Pt(32)
r.font.color.rgb = white

p = tf.add_paragraph()
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = f"Hari ini, Minggu {long_date}, ibadah akan dipimpin oleh {voorganger}." 
r.font.name = "Calibri"
r.font.size = Pt(32)
r.font.color.rgb = white

slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_background(slide)
box = slide.shapes.add_textbox(BOX_LEFT, Cm(1), BOX_WIDTH, Cm(16))
tf = box.text_frame
tf.clear()
tf.word_wrap = True
tf.auto_size = MSO_AUTO_SIZE.NONE
tf.vertical_anchor = MSO_ANCHOR.MIDDLE

p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = "Majelis yang bertugas dalam ibadah adalah "
r.font.name = "Calibri"
r.font.size = Pt(32)
r.font.color.rgb = white
r = p.add_run()
r.text = f"{ouderling}."
r.font.name = "Calibri"
r.font.size = Pt(32)
r.font.color.rgb = white
r.font.italic = True
r = p.add_run()
r.text = " Jika anda memiliki pertanyaan, dapat menghubungi beliau."
r.font.name = "Calibri"
r.font.size = Pt(32)
r.font.color.rgb = white

slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_background(slide)
box = slide.shapes.add_textbox(BOX_LEFT, Cm(1), BOX_WIDTH, Cm(16))
tf = box.text_frame
tf.clear()
tf.word_wrap = True
tf.auto_size = MSO_AUTO_SIZE.NONE
tf.vertical_anchor = MSO_ANCHOR.MIDDLE

p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = f"Hari Minggu, yang akan datang, {long_datenw}, ibadah akan dipimpin oleh ... Kebaktian akan dimulai pada pukul 10.30.\n"
r.font.name = "Calibri"
r.font.size = Pt(32)
r.font.color.rgb = white

p = tf.add_paragraph()
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = f"Pada hari minggu {long_datenw} yang akan datang, Online Landelijke Eredienst (OLE) akan dipimpin oleh …. dan dimulai pukul …"
r.font.name = "Calibri"
r.font.size = Pt(32)
r.font.color.rgb = white


# --- dankoffer vorige week ---

slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_background(slide)

box = slide.shapes.add_textbox(BOX_LEFT, Cm(0), BOX_WIDTH, Cm(4))
tf = box.text_frame
tf.clear()
tf.word_wrap = True
tf.auto_size = MSO_AUTO_SIZE.NONE
tf.vertical_anchor = MSO_ANCHOR.MIDDLE

p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = "Persembahan\nGKIN Amstelveen"
r.font.name = "Calibri"
r.font.size = Pt(36)
r.font.color.rgb = white
r.font.bold = True

rows, cols = 9, 3
table = slide.shapes.add_table(rows, cols, BOX_LEFT, Cm(4), BOX_WIDTH, Cm(10)).table
table.columns[0].width = Inches(4.6)
table.columns[1].width = Inches(0.7)
table.columns[2].width = Inches(1.0)

dark_blue = RGBColor(22, 25, 90)

lastweek = dienst_date + timedelta(days = -7)
long_datepw = format_date_long_nl(lastweek)

content = [
    [f"Minggu {long_datepw}", "", ""],
    ["1. Uang kontan", "€", "..."],
    ["2. Bon persembahan", "€", "..."],
    ["3. Transfer melalui bank", "€", "..."],
    ["4. TIKKIE", "€", "..."],
    ["Total Persembahan", "€", "..."],
    ["Jumlah jemaat yang hadir ... orang: ... dewasa, ... anak-anak.", "", ""],
    ["Penerimaan Gemeente Bijbelstudie ...", "€", "..."],
    [f"Penerimaan OLE {long_datepw}", "€", "..."]
]

# --- Title row (merge columns 0–2) ---
cell_start = table.cell(0, 0)
cell_end = table.cell(0, 2)
cell_start.merge(cell_end)  # merge done in-place
merged = cell_start          # keep reference to the surviving cell

merged.text = content[0][0]
p = merged.text_frame.paragraphs[0]
p.font.name = "Calibri"
p.font.bold = True
p.font.size = Pt(20)
p.font.color.rgb = white
merged.fill.background()

# --- Attendance row (merge columns 0–2) ---
cell_start = table.cell(6, 0)
cell_end = table.cell(6, 2)
cell_start.merge(cell_end)
merged = cell_start

merged.text = content[6][0]
p = merged.text_frame.paragraphs[0]
p.font.name = "Calibri"
p.font.size = Pt(16)
p.font.color.rgb = white
merged.fill.background()

# --- Handle the rest ---
for i, row_data in enumerate(content):
    if i in (0, 6):
        continue  # already merged and handled
    for j, text in enumerate(row_data):
        cell = table.cell(i, j)
        cell.text = text
        p = cell.text_frame.paragraphs[0]
        p.font.color.rgb = white
        p.font.name = "Calibri"
        cell.fill.background()

        # font size and style
        if i == 5:
            p.font.bold = True
            p.font.size = Pt(20)
        else:
            p.font.size = Pt(20)

        # alignment
        p.alignment = PP_ALIGN.RIGHT if j in (1, 2) else PP_ALIGN.LEFT


# --- dankoffer deze week ---

slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_background(slide)

box = slide.shapes.add_textbox(BOX_LEFT, Cm(0), BOX_WIDTH, Cm(4))
tf = box.text_frame
tf.clear()
tf.word_wrap = True
tf.auto_size = MSO_AUTO_SIZE.NONE
tf.vertical_anchor = MSO_ANCHOR.MIDDLE

p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = "Persembahan"
r.font.name = "Calibri"
r.font.size = Pt(36)
r.font.color.rgb = white
r.font.bold = True

box = slide.shapes.add_textbox(BOX_LEFT, Cm(3), BOX_WIDTH, Cm(16))
tf = box.text_frame
tf.clear()
tf.word_wrap = True
tf.auto_size = MSO_AUTO_SIZE.NONE
tf.vertical_anchor = MSO_ANCHOR.MIDDLE

p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = f"""Persembahan untuk pelayanan gereja dapat anda berikan dalam ibadah ini lewat 3 cara:
1.  Memasukkan dalam kantong persembahan yang ada di gereja.
2.	Melakukan transfer ke rekening GKIN Amstelveen,
IBAN: {IBAN}, 
Dengan mencantumkan tulisan “Collecte {short_date}”
"""
r.font.name = "Calibri"
r.font.size = Pt(32)
r.font.color.rgb = white

slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_background(slide)

box = slide.shapes.add_textbox(BOX_LEFT, Cm(3), BOX_WIDTH, Cm(16))
tf = box.text_frame
tf.clear()
tf.word_wrap = True
tf.auto_size = MSO_AUTO_SIZE.NONE
tf.vertical_anchor = MSO_ANCHOR.TOP

p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = f"""3.	Menggunakan QR-code atau tautan pembayaran:"""
r.font.name = "Calibri"
r.font.size = Pt(32)
r.font.color.rgb = white


# --- Regional ---

slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_background(slide)

box = slide.shapes.add_textbox(BOX_LEFT, Cm(0), BOX_WIDTH, Cm(17))
tf = box.text_frame
tf.clear()
tf.word_wrap = True
tf.auto_size = MSO_AUTO_SIZE.NONE
tf.vertical_anchor = MSO_ANCHOR.MIDDLE

p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = "PEMBERITAHUAN\nREGIONAL"
r.font.name = "Calibri"
r.font.size = Pt(48)
r.font.color.rgb = white
r.font.bold = True

# --- Overleden ---
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_background(slide)

box = slide.shapes.add_textbox(BOX_LEFT, Cm(0), BOX_WIDTH, Cm(2))
tf = box.text_frame
tf.clear()
tf.word_wrap = True
tf.auto_size = MSO_AUTO_SIZE.NONE
tf.vertical_anchor = MSO_ANCHOR.MIDDLE

p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = "Berita Duka"
r.font.name = "Calibri"
r.font.size = Pt(32)
r.font.color.rgb = white
r.font.bold = True

box = slide.shapes.add_textbox(BOX_LEFT, Cm(3), BOX_WIDTH, Cm(16))
tf = box.text_frame
tf.clear()
tf.word_wrap = True
tf.auto_size = MSO_AUTO_SIZE.NONE
tf.vertical_anchor = MSO_ANCHOR.MIDDLE

p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = """Telah meninggal dunia pada usia ... tahun, pada hari ..., ...: ..., dari regio Amstelveen, [isi hubungan keluarga disini].\n
Majelis dan Jemaat GKIN mengucapkan turut berduka cita yang sedalam-dalamnya. Kiranya Tuhan senantiasa menguatkan dan menghibur keluarga yang ditinggalkan.
"""
r.font.name = "Calibri"
r.font.size = Pt(32)
r.font.color.rgb = white

# --- Geboorte ---
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_background(slide)

box = slide.shapes.add_textbox(BOX_LEFT, Cm(0), BOX_WIDTH, Cm(2))
tf = box.text_frame
tf.clear()
tf.word_wrap = True
tf.auto_size = MSO_AUTO_SIZE.NONE
tf.vertical_anchor = MSO_ANCHOR.MIDDLE

p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = "Berita Kelahiran"
r.font.name = "Calibri"
r.font.size = Pt(32)
r.font.color.rgb = white
r.font.bold = True

box = slide.shapes.add_textbox(BOX_LEFT, Cm(3), BOX_WIDTH, Cm(16))
tf = box.text_frame
tf.clear()
tf.word_wrap = True
tf.auto_size = MSO_AUTO_SIZE.NONE
tf.vertical_anchor = MSO_ANCHOR.MIDDLE

p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = """Pada hari ..., ... telah lahir di ..., ..., putra dari ... dan ..., cucu dari..., dari GKIN regio Amstelveen.\n
Jemaat dan majelis gereja mengucapkan selamat kepada keluarga ..., kiranya Tuhan memberkati.
"""
r.font.name = "Calibri"
r.font.size = Pt(32)
r.font.color.rgb = white

# --- Overige ---
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_background(slide)

box = slide.shapes.add_textbox(BOX_LEFT, Cm(0), BOX_WIDTH, Cm(2))
tf = box.text_frame
tf.clear()
tf.word_wrap = True
tf.auto_size = MSO_AUTO_SIZE.NONE
tf.vertical_anchor = MSO_ANCHOR.MIDDLE

p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = "Titel"
r.font.name = "Calibri"
r.font.size = Pt(32)
r.font.color.rgb = white
r.font.bold = True

box = slide.shapes.add_textbox(BOX_LEFT, Cm(3), BOX_WIDTH, Cm(16))
tf = box.text_frame
tf.clear()
tf.word_wrap = True
tf.auto_size = MSO_AUTO_SIZE.NONE
tf.vertical_anchor = MSO_ANCHOR.MIDDLE

p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = """[Content... Lorem ipsum...]
"""
r.font.name = "Calibri"
r.font.size = Pt(32)
r.font.color.rgb = white

# --- Landelijk ---

slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_background(slide)

box = slide.shapes.add_textbox(BOX_LEFT, Cm(0), BOX_WIDTH, Cm(17))
tf = box.text_frame
tf.clear()
tf.word_wrap = True
tf.auto_size = MSO_AUTO_SIZE.NONE
tf.vertical_anchor = MSO_ANCHOR.MIDDLE

p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = "PEMBERITAHUAN\nLANDELIJK"
r.font.name = "Calibri"
r.font.size = Pt(48)
r.font.color.rgb = white
r.font.bold = True


# --- Overleden ---
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_background(slide)

box = slide.shapes.add_textbox(BOX_LEFT, Cm(0), BOX_WIDTH, Cm(2))
tf = box.text_frame
tf.clear()
tf.word_wrap = True
tf.auto_size = MSO_AUTO_SIZE.NONE
tf.vertical_anchor = MSO_ANCHOR.MIDDLE

p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = "Berita Duka"
r.font.name = "Calibri"
r.font.size = Pt(32)
r.font.color.rgb = white
r.font.bold = True

box = slide.shapes.add_textbox(BOX_LEFT, Cm(3), BOX_WIDTH, Cm(16))
tf = box.text_frame
tf.clear()
tf.word_wrap = True
tf.auto_size = MSO_AUTO_SIZE.NONE
tf.vertical_anchor = MSO_ANCHOR.MIDDLE

p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = """Telah meninggal dunia pada usia ... tahun, pada hari ..., ...: ..., dari regio ..., [isi hubungan keluarga disini].\n
Majelis dan Jemaat GKIN mengucapkan turut berduka cita yang sedalam-dalamnya. Kiranya Tuhan senantiasa menguatkan dan menghibur keluarga yang ditinggalkan.
"""
r.font.name = "Calibri"
r.font.size = Pt(32)
r.font.color.rgb = white

# --- Geboorte ---
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_background(slide)

box = slide.shapes.add_textbox(BOX_LEFT, Cm(0), BOX_WIDTH, Cm(2))
tf = box.text_frame
tf.clear()
tf.word_wrap = True
tf.auto_size = MSO_AUTO_SIZE.NONE
tf.vertical_anchor = MSO_ANCHOR.MIDDLE

p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = "Berita Kelahiran"
r.font.name = "Calibri"
r.font.size = Pt(32)
r.font.color.rgb = white
r.font.bold = True

box = slide.shapes.add_textbox(BOX_LEFT, Cm(3), BOX_WIDTH, Cm(16))
tf = box.text_frame
tf.clear()
tf.word_wrap = True
tf.auto_size = MSO_AUTO_SIZE.NONE
tf.vertical_anchor = MSO_ANCHOR.MIDDLE

p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = """Pada hari ..., ... telah lahir di ..., ..., putra dari ... dan ..., cucu dari..., dari GKIN regio ....\n
Jemaat dan majelis gereja mengucapkan selamat kepada keluarga ..., kiranya Tuhan memberkati.
"""
r.font.name = "Calibri"
r.font.size = Pt(32)
r.font.color.rgb = white

# --- Overige ---
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_background(slide)

box = slide.shapes.add_textbox(BOX_LEFT, Cm(0), BOX_WIDTH, Cm(2))
tf = box.text_frame
tf.clear()
tf.word_wrap = True
tf.auto_size = MSO_AUTO_SIZE.NONE
tf.vertical_anchor = MSO_ANCHOR.MIDDLE

p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = "Titel"
r.font.name = "Calibri"
r.font.size = Pt(32)
r.font.color.rgb = white
r.font.bold = True

box = slide.shapes.add_textbox(BOX_LEFT, Cm(3), BOX_WIDTH, Cm(16))
tf = box.text_frame
tf.clear()
tf.word_wrap = True
tf.auto_size = MSO_AUTO_SIZE.NONE
tf.vertical_anchor = MSO_ANCHOR.MIDDLE

p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = """[Content... Lorem ipsum...]
"""
r.font.name = "Calibri"
r.font.size = Pt(32)
r.font.color.rgb = white

# --- DIENST ---

def addBox(slide, text = "Text", left=Cm(0), top=Cm(0), width=Cm(0), height=Cm(0), color = RGBColor(255,255,255), Size = Pt(24), alignment = "CENTER", italic = False, bold = False, underline=False):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap =True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Size
    p.font.color.rgb = color
    p.font.italic = italic
    p.font.bold = bold
    p.font.underline = underline
    p.alignment = getattr(PP_ALIGN, alignment)
    p.font.name = "Calibri"

VERSE_OR_REFRRAIN_PREFIX = re.compile(
    r"""^\s*(                     # start, optional spaces
         (?:\d+\.)               # 1. 2. 3. ...
        |(?:[Rr]eff(?:rain)?[:.])# Reff:, Reff., reff:, reff., Reffrain:, reffrain.
       )\s*(.*)$                 # optional spaces then the rest of the line
    """,
    re.VERBOSE,
)

def split_prefix_to_own_line(line: str) -> str:
    """
    If a line starts with 'N.' or a refrain marker (Reff:, Reff., reff:, reff., Reffrain:, reffrain.),
    return 'prefix\\nrest'. Otherwise return the original line.
    """
    m = VERSE_OR_REFRRAIN_PREFIX.match(line)
    if not m:
        return line
    prefix, rest = m.group(1), m.group(2)
    # Even if rest is empty, keep the newline to put lyrics (if any) on the next line.
    return f"{prefix}\n{rest}" if rest is not None else f"{prefix}\n"


def add_song_slides(songNumber = 1, presentation = prs, staan = False):
    if not (1 <= songNumber <= len(songs)):
        raise ValueError(f"Song number {songNumber} out of range (have {len(songs)}).")
    s = songs[songNumber - 1]

    #Add slide for title
    slideT = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slideT)

    if s["Boek"] == "OPW":
        s["Boek"] = "Opwekking"

    if s["Boek"] == "Kosong":
        _m = re.search(r"\((.*?)\)", s["Titel"])
        if _m:
            s["Boek"] = _m.group(1)
            s["Titel"] = re.sub(r"\(.*?\)", "", s["Titel"]).strip()
        else:
            s["Boek"] = ""

        
    if (staan):
        addBox(slideT, "(gemeente gaat staan)", BOX_LEFT, Cm(3), BOX_WIDTH, Cm(2), RGBColor(255,173,3), italic=True, Size=Pt(24))
    
    if songNumber == 1:
        addBox(slideT, "INTOCHTSLIED", BOX_LEFT, Cm(4), BOX_WIDTH, Cm(2), bold=True, underline=True, Size=Pt(40))
    elif songNumber == 7:
        addBox(slideT, "SLOTLIED", BOX_LEFT, Cm(4), BOX_WIDTH, Cm(2), bold=True, underline=True, Size=Pt(40))
    else:
        addBox(slideT, "SAMENZANG", BOX_LEFT, Cm(4), BOX_WIDTH, Cm(2), bold=True, Size=Pt(40))

    if is_empty(s["Nummer"]):
        addBox(slideT, f'{s["Boek"]}\n\n "{s["Titel"]}"', BOX_LEFT, Cm(7), BOX_WIDTH, Cm(7), Size = Pt(40))
    else:
        if is_empty(s["Versen"]):
            addBox(slideT, f'{s["Boek"]} {s["Nummer"]}\n\n "{s["Titel"]}"', BOX_LEFT, Cm(7), BOX_WIDTH, Cm(7), Size = Pt(40))
        else:
            addBox(slideT, f'{s["Boek"]} {s["Nummer"]}: {s["Versen"]}\n\n  "{s["Titel"]}"', BOX_LEFT, Cm(7), BOX_WIDTH, Cm(7), Size = Pt(40))
    

    for si, song_text in enumerate(s["cells"]):
        blocks = [b.strip() for b in re.split(r"\r?\n\s*\r?\n", song_text.strip()) if b.strip()]
        
        for bi, block in enumerate(blocks):
            slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
            set_background(slide)

            if is_empty(s["Nummer"]):
                tt1 = f'{s["Boek"]}'
            else:
                tt1 = f'{s["Boek"]} {s["Nummer"]}'

            addBox(
                slide,
                text=tt1,
                left=Cm(17), top=Cm(0.5), width=Cm(6.7), height=Cm(2),
                alignment="RIGHT",
                color=RGBColor(255,173,3), bold=True, Size=Pt(24)
            )

            box = slide.shapes.add_textbox(BOX_LEFT, Cm(2), BOX_WIDTH, Cm(13.5))
            tf = box.text_frame
            tf.clear()
            tf.word_wrap = True
            tf.auto_size = MSO_AUTO_SIZE.NONE
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE

            lines = block.splitlines()
            for li, line in enumerate(lines):
                # NEW: normalize line so verse/refrain prefixes go on their own line
                line_out = split_prefix_to_own_line(line)

                p = tf.paragraphs[0] if li == 0 else tf.add_paragraph()
                p.alignment = PP_ALIGN.CENTER
                r = p.add_run()

                # keep your asterisk rule on the very last line of the last block
                if bi == len(blocks) - 1 and li == len(lines) - 1 and si == len(s["cells"]) - 1:
                    r.text = line_out + "\n\n*"
                else:
                    r.text = line_out

                r.font.name = "Calibri"
                r.font.size = Pt(32)
                r.font.color.rgb = white


# --- Intochtslied ---
add_song_slides(1, prs, True)

slide2 = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_background(slide2)

# --- Votum ---
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_background(slide)

box = slide.shapes.add_textbox(BOX_LEFT, Cm(1), BOX_WIDTH, Cm(2))
tf = box.text_frame
tf.clear()
tf.word_wrap = True
tf.auto_size = MSO_AUTO_SIZE.NONE
tf.vertical_anchor = MSO_ANCHOR.MIDDLE

p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = "(staande)"
r.font.name = "Calibri"
r.font.size = Pt(27)
r.font.color.rgb = yellow
r.font.italic = True

box = slide.shapes.add_textbox(BOX_LEFT, Cm(5), BOX_WIDTH, Cm(7))
tf = box.text_frame
tf.clear()
tf.word_wrap = True
tf.auto_size = MSO_AUTO_SIZE.NONE
tf.vertical_anchor = MSO_ANCHOR.MIDDLE

p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = "STIL GEBED\n"
r.font.name = "Calibri"
r.font.size = Pt(40)
r.font.color.rgb = white
r.font.bold = True

p = tf.add_paragraph()
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = "DOA HENING"
r.font.name = "Calibri"
r.font.size = Pt(40)
r.font.color.rgb = white
r.font.bold = True
r.font.italic = True

slide2 = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_background(slide2)

box = add_textbox(slide2, BOX_LEFT, Cm(2), BOX_WIDTH, Cm(1), "(staande)", "Calibri", 27, bold=False, alignment="CENTER", color=yellow)

box = slide2.shapes.add_textbox(BOX_LEFT, Cm(5), BOX_WIDTH, Cm(7))
tf = box.text_frame
tf.clear()
tf.word_wrap = True
tf.auto_size = MSO_AUTO_SIZE.NONE
tf.vertical_anchor = MSO_ANCHOR.MIDDLE

p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = "VOTUM EN GROET\n"
r.font.name = "Calibri"
r.font.size = Pt(40)
r.font.color.rgb = white
r.font.bold = True

p = tf.add_paragraph()
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = "VOTUM DAN SALAM"
r.font.name = "Calibri"
r.font.size = Pt(40)
r.font.color.rgb = white
r.font.bold = True
r.font.italic = True

p = tf.add_paragraph()
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = "("
r.font.name = "Calibri"
r.font.size = Pt(32)
r.font.color.rgb = white
r.font.italic = True

r = p.add_run()
r.text = "V: Voorganger"
r.font.name = "Calibri"
r.font.size = Pt(32)
r.font.color.rgb = yellow
r.font.italic = True

r = p.add_run()
r.text = ", G: Gemeente)"
r.font.name = "Calibri"
r.font.size = Pt(32)
r.font.color.rgb = white
r.font.italic = True

slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_background(slide)

box = slide.shapes.add_textbox(BOX_LEFT, Cm(1), BOX_WIDTH, Cm(16))
tf = box.text_frame
tf.clear()
tf.word_wrap = True
tf.auto_size = MSO_AUTO_SIZE.NONE
tf.vertical_anchor = MSO_ANCHOR.MIDDLE

p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = "V: Onze hulp is in de Naam van de Heer,\n"
r.font.name = "Calibri"
r.font.size = Pt(32)
r.font.color.rgb = yellow
r = p.add_run()
r.text = "(Pertolongan kita adalah dalam nama Tuhan)\n\n"
r.font.name = "Calibri"
r.font.size = Pt(32)
r.font.color.rgb = yellow
r.font.italic = True
r = p.add_run()
r.text = "G: die hemel en aarde gemaakt heeft,\n"
r.font.name = "Calibri"
r.font.size = Pt(32)
r.font.color.rgb = white
r.font.bold = True
r = p.add_run()
r.text = "(yang menjadikan langit dan bumi)"
r.font.name = "Calibri"
r.font.size = Pt(32)
r.font.color.rgb = white
r.font.bold = True
r.font.italic = True

slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_background(slide)

box = slide.shapes.add_textbox(BOX_LEFT, Cm(1), BOX_WIDTH, Cm(16))
tf = box.text_frame
tf.clear()
tf.word_wrap = True
tf.auto_size = MSO_AUTO_SIZE.NONE
tf.vertical_anchor = MSO_ANCHOR.MIDDLE

p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = "V: die trouw houdt tot in eeuwigheid,\n"
r.font.name = "Calibri"
r.font.size = Pt(32)
r.font.color.rgb = yellow
r = p.add_run()
r.text = "(yang setia dan tetap setia sampai kekal selamanya)\n\n"
r.font.name = "Calibri"
r.font.size = Pt(32)
r.font.color.rgb = yellow
r.font.italic = True
r = p.add_run()
r.text = "G: en niet laat varen het werk van Zijn handen,\n"
r.font.name = "Calibri"
r.font.size = Pt(32)
r.font.color.rgb = white
r.font.bold = True
r = p.add_run()
r.text = "(dan yang tidak meninggalkan perbuatan tanganNya)"
r.font.name = "Calibri"
r.font.size = Pt(32)
r.font.color.rgb = white
r.font.bold = True
r.font.italic = True

slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_background(slide)

box = slide.shapes.add_textbox(BOX_LEFT, Cm(1), BOX_WIDTH, Cm(16))
tf = box.text_frame
tf.clear()
tf.word_wrap = True
tf.auto_size = MSO_AUTO_SIZE.NONE
tf.vertical_anchor = MSO_ANCHOR.MIDDLE

p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = "V: Genade zij u en vrede van God, onze Vader, en van de Heer Jezus Christus, in de gemeenschap van de Heilige Geest.\n"
r.font.name = "Calibri"
r.font.size = Pt(32)
r.font.color.rgb = yellow
r = p.add_run()
r.text = "(Kasih karunia dan damai sejahtera dari Allah Bapa kita dan dari Tuhan kita Yesus Kristus, dalam persekutuan Roh Kudus menyertai saudara sekalian)\n\n"
r.font.name = "Calibri"
r.font.size = Pt(32)
r.font.color.rgb = yellow
r.font.italic = True
r = p.add_run()
r.text = "G: Amen"
r.font.name = "Calibri"
r.font.size = Pt(32)
r.font.color.rgb = white
r.font.bold = True

# --- Aanvangstekst ---
boek = _get(table2, 0, ["Boek", "boek"])
hs   = _get(table2, 0, ["H.s.", "H.s", "HS", "H.S.", "H.S"])
vfrom = _get(table2, 0, ["Vers van", "Vers Van", "vers van", "Van"])
vto   = _get(table2, 0, ["Vers tot", "Vers Tot", "vers tot", "Tot"])
boek2 = _get(table2, 0, ["Boek 2", "boek 2"])
hs2   = _get(table2, 0, ["H.s. 2", "H.s 2", "HS 2", "H.S. 2", "H.S 2"])
vfrom2 = _get(table2, 0, ["Vers van 2", "Vers Van 2", "vers van 2", "Van 2"])
vto2   = _get(table2, 0, ["Vers tot 2", "Vers Tot 2", "vers tot 2", "Tot 2"])
boek_id = dutch_to_indonesian_bible.get(boek, boek)
boek_id2 = dutch_to_indonesian_bible.get(boek2, boek2)

slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_background(slide)

box = slide.shapes.add_textbox(BOX_LEFT, Cm(1), BOX_WIDTH, Cm(16))
tf = box.text_frame
tf.clear()
tf.word_wrap = True
tf.auto_size = MSO_AUTO_SIZE.NONE
tf.vertical_anchor = MSO_ANCHOR.TOP
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = "(zittende)\n"
r.font.name = "Calibri"
r.font.size = Pt(28)
r.font.color.rgb = yellow
r.font.italic = True

p = tf.add_paragraph()
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = "\nAANVANGSTEKST\n\n"
r.font.name = "Calibri"
r.font.size = Pt(40)
r.font.color.rgb = white
r.font.bold = True
r = p.add_run()
r.text = "AYAT PEMBUKA\n"
r.font.name = "Calibri"
r.font.size = Pt(40)
r.font.color.rgb = white
r.font.bold = True
r.font.italic = True

text = ""

if vto and str(vto).strip():
    text = f"{boek} / {boek_id} {hs}: {vfrom}-{vto}"
else:
    text = f"{boek} / {boek_id} {hs}: {vfrom}"

if is_empty(boek2) == False:
    if vto2 and str(vto2).strip():
        text = text + f"\nen {boek2} / {boek_id2} {hs2}: {vfrom2}-{vto2}"
    else:
        text = text + f"\nen {boek2} / {boek_id2} {hs2}: {vfrom2}"

p = tf.add_paragraph()
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = f"\n{text}"
r.font.name = "Calibri"
r.font.size = Pt(40)
r.font.color.rgb = white


add_verses_to_ppt(prs, boek, hs, vfrom, vto)
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_background(slide)
add_verses_to_ppt_indo(prs, boek_id, hs, vfrom, vto)
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_background(slide)

if is_empty(boek2) == False:
    add_verses_to_ppt(prs, boek2, hs2, vfrom2, vto2)
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_background(slide)
    add_verses_to_ppt_indo(prs, boek_id2, hs2, vfrom2, vto2)
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_background(slide)


# --- 2e lied ---
add_song_slides(2, prs, False)

slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_background(slide)


# --- Gebed van Toenadering ---

slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_background(slide)

box = slide.shapes.add_textbox(BOX_LEFT, Cm(1), BOX_WIDTH, Cm(16))
tf = box.text_frame
tf.clear()
tf.word_wrap = True
tf.auto_size = MSO_AUTO_SIZE.NONE
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run()

if is_empty(tieners):
    r.text = "De kinderen gaan naar de kindernevendienst\n\n Er is geen tienerdienst"
else:
    r.text = "De kinderen gaan naar de kindernevendienst\n\n Tieners gaan naar de tienerdienst"

r.font.name = "Calibri"
r.font.size = Pt(40)
r.font.color.rgb = white
r.font.bold = True

slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_background(slide)

box = slide.shapes.add_textbox(BOX_LEFT, Cm(1), BOX_WIDTH, Cm(16))
tf = box.text_frame
tf.clear()
tf.word_wrap = True
tf.auto_size = MSO_AUTO_SIZE.NONE
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = "GEBED VAN TOENADERING"
r.font.name = "Calibri"
r.font.size = Pt(40)
r.font.color.rgb = white
r.font.bold = True
r = p.add_run()
r.text = "\n\nDOA MENGHADAP TUHAN"
r.font.name = "Calibri"
r.font.size = Pt(40)
r.font.color.rgb = white
r.font.bold = True
r.font.italic = True

slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_background(slide)

# --- 3e lied ---

add_song_slides(3, prs, False)

slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_background(slide)

# --- GENADEVERKONDIGING ---
boek = _get(table2, 1, ["Boek", "boek"])
hs   = _get(table2, 1, ["H.s.", "H.s", "HS", "H.S.", "H.S"])
vfrom = _get(table2, 1, ["Vers van", "Vers Van", "vers van", "Van"])
vto   = _get(table2, 1, ["Vers tot", "Vers Tot", "vers tot", "Tot"])
boek2 = _get(table2, 1, ["Boek 2", "boek 2"])
hs2   = _get(table2, 1, ["H.s. 2", "H.s 2", "HS 2", "H.S. 2", "H.S 2"])
vfrom2 = _get(table2, 1, ["Vers van 2", "Vers Van 2", "vers van 2", "Van 2"])
vto2   = _get(table2, 1, ["Vers tot 2", "Vers Tot 2", "vers tot 2", "Tot 2"])
boek_id = dutch_to_indonesian_bible.get(boek, boek)
boek_id2 = dutch_to_indonesian_bible.get(boek2, boek2)

slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_background(slide)

box = slide.shapes.add_textbox(BOX_LEFT, Cm(1), BOX_WIDTH, Cm(16))
tf = box.text_frame
tf.clear()
tf.word_wrap = True
tf.auto_size = MSO_AUTO_SIZE.NONE
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = "GENADEVERKONDIGING\n"
r.font.name = "Calibri"
r.font.size = Pt(40)
r.font.color.rgb = white
r.font.bold = True
r = p.add_run()
r.text = "\nBERITA ANUGERAH\n"
r.font.name = "Calibri"
r.font.size = Pt(40)
r.font.color.rgb = white
r.font.bold = True
r.font.italic = True

text = ""

if vto and str(vto).strip():
    text = f"{boek} / {boek_id} {hs}: {vfrom}-{vto}"
else:
    text = f"{boek} / {boek_id} {hs}: {vfrom}"

if is_empty(boek2) == False:
    if vto2 and str(vto2).strip():
        text = text + f"\nen {boek2} / {boek_id2} {hs2}: {vfrom2}-{vto2}"
    else:
        text = text + f"\nen {boek2} / {boek_id2} {hs2}: {vfrom2}"

p = tf.add_paragraph()
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = f"\n{text}"
r.font.name = "Calibri"
r.font.size = Pt(40)
r.font.color.rgb = white

add_verses_to_ppt(prs, boek, hs, vfrom, vto)
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_background(slide)
add_verses_to_ppt_indo(prs, boek_id, hs, vfrom, vto)
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_background(slide)

if is_empty(boek2) == False:
    add_verses_to_ppt(prs, boek2, hs2, vfrom2, vto2)
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_background(slide)
    add_verses_to_ppt_indo(prs, boek_id2, hs2, vfrom2, vto2)
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_background(slide)

# --- Levensvernieuwing ---
boek = _get(table2, 2, ["Boek", "boek"])
hs   = _get(table2, 2, ["H.s.", "H.s", "HS", "H.S.", "H.S"])
vfrom = _get(table2, 2, ["Vers van", "Vers Van", "vers van", "Van"])
vto   = _get(table2, 2, ["Vers tot", "Vers Tot", "vers tot", "Tot"])
boek2 = _get(table2, 2, ["Boek 2", "boek 2"])
hs2   = _get(table2, 2, ["H.s. 2", "H.s 2", "HS 2", "H.S. 2", "H.S 2"])
vfrom2 = _get(table2, 2, ["Vers van 2", "Vers Van 2", "vers van 2", "Van 2"])
vto2   = _get(table2, 2, ["Vers tot 2", "Vers Tot 2", "vers tot 2", "Tot 2"])
boek_id = dutch_to_indonesian_bible.get(boek, boek)
boek_id2 = dutch_to_indonesian_bible.get(boek2, boek2)

slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_background(slide)

box = slide.shapes.add_textbox(BOX_LEFT, Cm(1), BOX_WIDTH, Cm(16))
tf = box.text_frame
tf.clear()
tf.word_wrap = True
tf.auto_size = MSO_AUTO_SIZE.NONE
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = "GODS WOORD VOOR LEVENSVERNIEUWING\n"
r.font.name = "Calibri"
r.font.size = Pt(40)
r.font.color.rgb = white
r.font.bold = True

r = p.add_run()
r.text = "\nPETUNJUK HIDUP BARU\n"
r.font.name = "Calibri"
r.font.size = Pt(40)
r.font.color.rgb = white
r.font.bold = True
r.font.italic = True

text = ""

if vto and str(vto).strip():
    text = f"{boek} / {boek_id} {hs}: {vfrom}-{vto}"
else:
    text = f"{boek} / {boek_id} {hs}: {vfrom}"

if is_empty(boek2) == False:
    if vto2 and str(vto2).strip():
        text = text + f"\nen {boek2} / {boek_id2} {hs2}: {vfrom2}-{vto2}"
    else:
        text = text + f"\nen {boek2} / {boek_id2} {hs2}: {vfrom2}"

p = tf.add_paragraph()
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = f"\n{text}"
r.font.name = "Calibri"
r.font.size = Pt(40)
r.font.color.rgb = white

add_verses_to_ppt(prs, boek, hs, vfrom, vto)
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_background(slide)
add_verses_to_ppt_indo(prs, boek_id, hs, vfrom, vto)
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_background(slide)

if is_empty(boek2) == False:
    add_verses_to_ppt(prs, boek2, hs2, vfrom2, vto2)
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_background(slide)
    add_verses_to_ppt_indo(prs, boek_id2, hs2, vfrom2, vto2)
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_background(slide)

# --- 4e Lied ---
add_song_slides(4, prs, False)
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_background(slide)

# --- Schriftlezing ---
boek = _get(table2, 3, ["Boek", "boek"])
hs   = _get(table2, 3, ["H.s.", "H.s", "HS", "H.S.", "H.S"])
vfrom = _get(table2, 3, ["Vers van", "Vers Van", "vers van", "Van"])
vto   = _get(table2, 3, ["Vers tot", "Vers Tot", "vers tot", "Tot"])
boek2 = _get(table2, 3, ["Boek 2", "boek 2"])
hs2   = _get(table2, 3, ["H.s. 2", "H.s 2", "HS 2", "H.S. 2", "H.S 2"])
vfrom2 = _get(table2, 3, ["Vers van 2", "Vers Van 2", "vers van 2", "Van 2"])
vto2   = _get(table2, 3, ["Vers tot 2", "Vers Tot 2", "vers tot 2", "Tot 2"])
boek_id = dutch_to_indonesian_bible.get(boek, boek)
boek_id2 = dutch_to_indonesian_bible.get(boek2, boek2)

slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_background(slide)
box = slide.shapes.add_textbox(BOX_LEFT, Cm(1), BOX_WIDTH, Cm(16))
tf = box.text_frame
tf.clear()
tf.word_wrap = True
tf.auto_size = MSO_AUTO_SIZE.NONE
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = "GEBED OM VERLICHTING MET DE HEILIGE GEEST\n"
r.font.name = "Calibri"
r.font.size = Pt(40)
r.font.color.rgb = white
r.font.bold = True
r = p.add_run()
r.text = "\nDOA MOHON BIMBINGAN ROH KUDUS"
r.font.name = "Calibri"
r.font.size = Pt(40)
r.font.color.rgb = white
r.font.italic = True
r.font.bold = True

slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_background(slide)
box = slide.shapes.add_textbox(BOX_LEFT, Cm(1), BOX_WIDTH, Cm(16))
tf = box.text_frame
tf.clear()
tf.word_wrap = True
tf.auto_size = MSO_AUTO_SIZE.NONE
tf.vertical_anchor = MSO_ANCHOR.TOP
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = "SCHRIFTLEZING\n"
r.font.name = "Calibri"
r.font.size = Pt(40)
r.font.color.rgb = white
r.font.bold = True
r = p.add_run()
r.text = "\nPEMBACAAN ALKITAB\n"
r.font.name = "Calibri"
r.font.size = Pt(40)
r.font.color.rgb = white
r.font.italic = True
r.font.bold = True

text = ""

if vto and str(vto).strip():
    textP = f"{boek} / {boek_id} {hs}: {vfrom}-{vto}"
else:
    textP = f"{boek} / {boek_id} {hs}: {vfrom}"

if is_empty(boek2) == False:
    if vto2 and str(vto2).strip():
        textP = text + f"\nen {boek2} / {boek_id2} {hs2}: {vfrom2}-{vto2}"
    else:
        textP = text + f"\nen {boek2} / {boek_id2} {hs2}: {vfrom2}"

r = p.add_run()
r.text = textP
r.font.name = "Calibri"
r.font.size = Pt(36)
r.font.color.rgb = white

add_verses_to_ppt(prs, boek, hs, vfrom, vto)
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_background(slide)
add_verses_to_ppt_indo(prs, boek_id, hs, vfrom, vto)

if is_empty(boek2) == False:
    add_verses_to_ppt(prs, boek2, hs2, vfrom2, vto2)
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_background(slide)
    add_verses_to_ppt_indo(prs, boek_id2, hs2, vfrom2, vto2)
    
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_background(slide)
box = slide.shapes.add_textbox(BOX_LEFT, Cm(1), BOX_WIDTH, Cm(16))
tf = box.text_frame
tf.clear()
tf.word_wrap = True
tf.auto_size = MSO_AUTO_SIZE.NONE
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = "V: Gelukkig zij die luisteren naar het Woord van God en ernaar leven.\n"
r.font.name = "Calibri"
r.font.size = Pt(32)
r.font.color.rgb = yellow
r = p.add_run()
r.text = f"G. zingt: '{HMH} (3x)' uit KJ 473b"
r.font.name = "Calibri"
r.font.size = Pt(32)
r.font.color.rgb = white
r.font.italic = True

# --- Preek ---

slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_background(slide)
box = slide.shapes.add_textbox(BOX_LEFT, Cm(1), BOX_WIDTH, Cm(3))
tf = box.text_frame
tf.clear()
tf.word_wrap = True
tf.auto_size = MSO_AUTO_SIZE.NONE
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = "PREEK / KOTBAH"
r.font.name = "Calibri"
r.font.size = Pt(40)
r.font.color.rgb = yellow

box = slide.shapes.add_textbox(BOX_LEFT, Cm(3), BOX_WIDTH, Cm(13))
tf = box.text_frame
tf.clear()
tf.word_wrap = True
tf.auto_size = MSO_AUTO_SIZE.NONE
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = "Thema / Tema:\n\n"
r.font.name = "Calibri"
r.font.size = Pt(40)
r.font.color.rgb = white
r.font.bold = True
r = p.add_run()
r.text = f'{titelPred}\n\n{textP}'
r.font.name = "Calibri"
r.font.size = Pt(32)
r.font.color.rgb = white
r.font.italic = True



path_sermon = os.path.join(dir_path + '/file mingguan/', PREEK_NAME)
if os.path.exists(path_sermon):
    add_sermon_doc_to_ppt(prs, path_sermon, box_width=BOX_WIDTH, box_left=BOX_LEFT, box_top=Cm(0.5), font_size_pt=28, max_lines_per_slide=16)


# --- Meditatief moment ---
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_background(slide)
box = slide.shapes.add_textbox(BOX_LEFT, Cm(1), BOX_WIDTH, Cm(16))
tf = box.text_frame
tf.clear()
tf.word_wrap = True
tf.auto_size = MSO_AUTO_SIZE.NONE
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = "MEDITATIEF MOMENT\n / Voordracht"
r.font.name = "Calibri"
r.font.size = Pt(32)
r.font.color.rgb = white
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_background(slide)

# --- Apostolische Geloofsbelijdenis
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_background(slide)
addBox(slide, "(staande)", BOX_LEFT, Cm(3), BOX_WIDTH, Cm(2), RGBColor(255,173,3), italic=True, Size=Pt(32))
box = slide.shapes.add_textbox(BOX_LEFT, Cm(1), BOX_WIDTH, Cm(16))
tf = box.text_frame
tf.clear()
tf.word_wrap = True
tf.auto_size = MSO_AUTO_SIZE.NONE
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = "APOSTOLISCHE\nGELOOFSBELIJDENIS\n"
r.font.name = "Calibri"
r.font.size = Pt(40)
r.font.color.rgb = white
r.font.bold = True
r = p.add_run()
r.text = "\nPENGAKUAN IMAN RASULI"
r.font.name = "Calibri"
r.font.size = Pt(40)
r.font.color.rgb = white
r.font.italic = True
r.font.bold = True

slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_background(slide)
box = slide.shapes.add_textbox(BOX_LEFT, Cm(1), BOX_WIDTH, Cm(16))
tf = box.text_frame
tf.clear()
tf.word_wrap = True
tf.auto_size = MSO_AUTO_SIZE.NONE
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = "Ik geloof  in God de Vader, de Almachtige, Schepper van de hemel en aarde.\nEn in Jezus Christus, zijn eniggeboren Zoon, onze Here; die ontvangen is van de Heilige Geest, geboren uit de maagd Maria;\ndie geleden heeft onder Pontius Pilatus,\nis gekruisigd, gestorven en begraven, nedergedaald naar het dodenrijk;"
r.font.name = "Calibri"
r.font.size = Pt(32)
r.font.color.rgb = white

slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_background(slide)
box = slide.shapes.add_textbox(BOX_LEFT, Cm(1), BOX_WIDTH, Cm(16))
tf = box.text_frame
tf.clear()
tf.word_wrap = True
tf.auto_size = MSO_AUTO_SIZE.NONE
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = "Op de derde dag opgestaan uit de doden;\nOpgevaren naar de hemel en zit aan de rechterhand van God, de almachtige Vader;\nVandaar zal Hij komen om te oordelen de levenden en de doden."
r.font.name = "Calibri"
r.font.size = Pt(32)
r.font.color.rgb = white

slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_background(slide)
box = slide.shapes.add_textbox(BOX_LEFT, Cm(1), BOX_WIDTH, Cm(16))
tf = box.text_frame
tf.clear()
tf.word_wrap = True
tf.auto_size = MSO_AUTO_SIZE.NONE
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = "Ik geloof  in de Heilige Geest;\n één heilige, algemene, christelijke kerk,\n de gemeenschap der heiligen; vergeving van de zonden;\nopstanding van het vlees; en een eeuwig leven.\n Amen"
r.font.name = "Calibri"
r.font.size = Pt(32)
r.font.color.rgb = white

slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_background(slide)

# --- 5e lied ---
add_song_slides(5, prs, False)
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_background(slide)

# --- Voorbeden ---
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_background(slide)
box = slide.shapes.add_textbox(BOX_LEFT, Cm(1), BOX_WIDTH, Cm(16))
tf = box.text_frame
tf.clear()
tf.word_wrap = True
tf.auto_size = MSO_AUTO_SIZE.NONE
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = "VOORBEDEN\n"
r.font.name = "Calibri"
r.font.size = Pt(40)
r.font.color.rgb = white
r.font.bold = True
r = p.add_run()
r.text = "\nDOA SYAFAAT\n\n"
r.font.name = "Calibri"
r.font.size = Pt(40)
r.font.color.rgb = white
r.font.bold = True
r.font.italic = True
r = p.add_run()
r.text = "(eindigend met het Onze Vader)\n"
r.font.name = "Calibri"
r.font.size = Pt(34)
r.font.color.rgb = white
r = p.add_run()
r.text = "(dan diakhiri doa Bapa Kami)"
r.font.name = "Calibri"
r.font.size = Pt(34)
r.font.color.rgb = white
r.font.italic = True


# --- Dankoffer ---

slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_background(slide)
box = slide.shapes.add_textbox(BOX_LEFT, Cm(1), BOX_WIDTH, Cm(16))
tf = box.text_frame
tf.clear()
tf.word_wrap = True
tf.auto_size = MSO_AUTO_SIZE.NONE
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = "DANKOFFER\n\n"
r.font.name = "Calibri"
r.font.size = Pt(40)
r.font.color.rgb = white
r.font.bold = True
r = p.add_run()
r.text = "\nPERSEMBAHAN\n\n"
r.font.name = "Calibri"
r.font.size = Pt(40)
r.font.color.rgb = white
r.font.bold = True
r.font.italic = True

boek = _get(table2, 4, ["Boek", "boek"])
hs   = _get(table2, 4, ["H.s.", "H.s", "HS", "H.S.", "H.S"])
vfrom = _get(table2, 4, ["Vers van", "Vers Van", "vers van", "Van"])
vto   = _get(table2, 4, ["Vers tot", "Vers Tot", "vers tot", "Tot"])

slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_background(slide)
box = slide.shapes.add_textbox(BOX_LEFT, Cm(1), BOX_WIDTH, Cm(16))
tf = box.text_frame
tf.clear()
tf.word_wrap = True
tf.auto_size = MSO_AUTO_SIZE.NONE
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run()

text=""

if vto and str(vto).strip():
    text = f"{boek} {hs}: {vfrom}-{vto}"
else:
    text = f"{boek} {hs}: {vfrom}"

r.text = f"Ouderling: Gemeente, wij krijgen nu de gelegenheid ons dankoffer aan God te brengen. En we gedenken daarbij de woorden uit {text}."
r.font.name = "Calibri"
r.font.size = Pt(32)
r.font.color.rgb = white

add_verses_to_ppt(prs, boek, hs, vfrom, vto)

slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_background(slide)
box = slide.shapes.add_textbox(BOX_LEFT, Cm(1), BOX_WIDTH, Cm(16))
tf = box.text_frame
tf.clear()
tf.word_wrap = True
tf.auto_size = MSO_AUTO_SIZE.NONE
tf.vertical_anchor = MSO_ANCHOR.TOP
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = f"Het dankoffer kunt u overmaken naar de rekeningnummer van Gereja Kristen Indonesia Nederland.\n IBAN: {IBAN}\no.v.v. Collecte {short_date}"
r.font.name = "Calibri"
r.font.size = Pt(32)
r.font.color.rgb = white

add_song_slides(6, prs, False)

slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_background(slide)
box = slide.shapes.add_textbox(BOX_LEFT, Cm(1), BOX_WIDTH, Cm(16))
tf = box.text_frame
tf.clear()
tf.word_wrap = True
tf.auto_size = MSO_AUTO_SIZE.NONE
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = "DANKGEBED\n"
r.font.name = "Calibri"
r.font.size = Pt(40)
r.font.color.rgb = white
r.font.bold = True

r = p.add_run()
r.text = "\nDOA PERSEMBAHAN\n\n"
r.font.name = "Calibri"
r.font.size = Pt(40)
r.font.color.rgb = white
r.font.bold = True
r.font.italic = True

slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_background(slide)

# --- 7e / Slot-lied ---
add_song_slides(7, prs, True)

slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_background(slide)

# --- Zending en Zegen ---
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_background(slide)
box = slide.shapes.add_textbox(BOX_LEFT, Cm(0), BOX_WIDTH, Cm(5))
tf = box.text_frame
tf.clear()
tf.word_wrap = True
tf.auto_size = MSO_AUTO_SIZE.NONE
tf.vertical_anchor = MSO_ANCHOR.TOP
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = "ZENDING & ZEGEN\n"
r.font.name = "Calibri"
r.font.size = Pt(40)
r.font.color.rgb = white
r.font.bold = True
r = p.add_run()
r.text = "PENGUTUSAN & BERKAT"
r.font.name = "Calibri"
r.font.size = Pt(40)
r.font.color.rgb = white
r.font.bold = True
r.font.italic = True

box = slide.shapes.add_textbox(BOX_LEFT, Cm(2), BOX_WIDTH, Cm(16))
tf = box.text_frame
tf.clear()
tf.word_wrap = True
tf.auto_size = MSO_AUTO_SIZE.NONE
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = "V: Verheft uw harten tot de Heer"
r.font.name = "Calibri"
r.font.size = Pt(28)
r.font.color.rgb = yellow
r=p.add_run()
r.text = "\n(Arahkan hatimu kepada Tuhan!)\n\n"
r.font.name = "Calibri"
r.font.size = Pt(28)
r.font.color.rgb = yellow
r.font.italic = True
r = p.add_run()
r.text = "G: Wij zijn met ons hart bij de Heer"
r.font.name = "Calibri"
r.font.size = Pt(28)
r.font.color.rgb = white
r.font.bold = True
r=p.add_run()
r.text = "\n(Kami mengarahkan hati kami kepada Tuhan)"
r.font.name = "Calibri"
r.font.size = Pt(28)
r.font.color.rgb = white
r.font.bold = True
r.font.italic = True

slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_background(slide)
box = slide.shapes.add_textbox(BOX_LEFT, Cm(0), BOX_WIDTH, Cm(5))
tf = box.text_frame
tf.clear()
tf.word_wrap = True
tf.auto_size = MSO_AUTO_SIZE.NONE
tf.vertical_anchor = MSO_ANCHOR.TOP
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = "ZENDING & ZEGEN\n"
r.font.name = "Calibri"
r.font.size = Pt(40)
r.font.color.rgb = white
r.font.bold = True
r = p.add_run()
r.text = "PENGUTUSAN & BERKAT"
r.font.name = "Calibri"
r.font.size = Pt(40)
r.font.color.rgb = white
r.font.bold = True
r.font.italic = True

box = slide.shapes.add_textbox(BOX_LEFT, Cm(2), BOX_WIDTH, Cm(16))
tf = box.text_frame
tf.clear()
tf.word_wrap = True
tf.auto_size = MSO_AUTO_SIZE.NONE
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = "V: Wees getuigen van Christus"
r.font.name = "Calibri"
r.font.size = Pt(24)
r.font.color.rgb = yellow
r=p.add_run()
r.text = "\n(Jadilah saksi Kristus)\n\n"
r.font.name = "Calibri"
r.font.size = Pt(24)
r.font.color.rgb = yellow
r.font.italic = True
r = p.add_run()
r.text = "G: Lof aan God"
r.font.name = "Calibri"
r.font.size = Pt(24)
r.font.color.rgb = white
r.font.bold = True
r=p.add_run()
r.text = "\n(Syukur kepada Allah)\n\n"
r.font.name = "Calibri"
r.font.size = Pt(24)
r.font.color.rgb = white
r.font.bold = True
r.font.italic = True
r = p.add_run()
r.text = "V: Geloofd zij de Heer"
r.font.name = "Calibri"
r.font.size = Pt(24)
r.font.color.rgb = yellow
r=p.add_run()
r.text = "\n(Terpujilah Tuhan)\n\n"
r.font.name = "Calibri"
r.font.size = Pt(24)
r.font.color.rgb = yellow
r.font.italic = True
r = p.add_run()
r.text = "G: Nu en voor altijd"
r.font.name = "Calibri"
r.font.size = Pt(24)
r.font.color.rgb = white
r.font.bold = True
r=p.add_run()
r.text = "\n(Kini dan selamanya)"
r.font.name = "Calibri"
r.font.size = Pt(24)
r.font.color.rgb = white
r.font.bold = True
r.font.italic = True
    
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_background(slide)
box = slide.shapes.add_textbox(BOX_LEFT, Cm(0), BOX_WIDTH, Cm(5))
tf = box.text_frame
tf.clear()
tf.word_wrap = True
tf.auto_size = MSO_AUTO_SIZE.NONE
tf.vertical_anchor = MSO_ANCHOR.TOP
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = "ZENDING & ZEGEN\n"
r.font.name = "Calibri"
r.font.size = Pt(40)
r.font.color.rgb = white
r.font.bold = True
r = p.add_run()
r.text = "PENGUTUSAN & BERKAT"
r.font.name = "Calibri"
r.font.size = Pt(40)
r.font.color.rgb = white
r.font.bold = True
r.font.italic = True

box = slide.shapes.add_textbox(BOX_LEFT, Cm(2), BOX_WIDTH, Cm(16))
tf = box.text_frame
tf.clear()
tf.word_wrap = True
tf.auto_size = MSO_AUTO_SIZE.NONE
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = "V: Ontvang nu de zegen van de Heer:"
r.font.name = "Calibri"
r.font.size = Pt(28)
r.font.color.rgb = yellow
r=p.add_run()
r.text = "\n(Terimalah berkat Tuhan):"
r.font.name = "Calibri"
r.font.size = Pt(28)
r.font.color.rgb = yellow
r.font.italic = True


slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_background(slide)

box = slide.shapes.add_textbox(BOX_LEFT, Cm(1), BOX_WIDTH, Cm(16))
tf = box.text_frame
tf.clear()
tf.word_wrap = True
tf.auto_size = MSO_AUTO_SIZE.NONE
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = "GEMEENTE ZINGT:\n"
r.font.name = "Calibri"
r.font.size = Pt(40)
r.font.color.rgb = white
r.font.bold = True
r = p.add_run()
r.text = '"AMEN, AMEN, AMEN"'
r.font.name = "Calibri"
r.font.size = Pt(40)
r.font.color.rgb = white
r.font.bold = True
r.font.italic = True

slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_background(slide)

box = slide.shapes.add_textbox(BOX_LEFT, Cm(1), BOX_WIDTH, Cm(16))
tf = box.text_frame
tf.clear()
tf.word_wrap = True
tf.auto_size = MSO_AUTO_SIZE.NONE
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = "STIL GEBED\n"
r.font.name = "Calibri"
r.font.size = Pt(40)
r.font.color.rgb = white
r.font.bold = True
r=p.add_run()
r.text = "\nDOA HENING"
r.font.name = "Calibri"
r.font.size = Pt(40)
r.font.color.rgb = white
r.font.italic = True
r.font.bold = True
r=p.add_run()
r.text = "\n(GEMEENTE BLIJFT STAAN EN HOUDT EEN MOMENT VAN STILTE)"
r.font.name = "Calibri"
r.font.size = Pt(40)
r.font.color.rgb = white
r.font.bold = True


# Save after you’ve added all slides
prs_name = f"LiturgieP {short_date}_{REGIO}.pptx"
out_path = os.path.join(dir_path + '/file mingguan/', prs_name)
prs.save(out_path)
print(f"Powerpoint presentation generated: {out_path}")
